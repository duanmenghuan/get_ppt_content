from pptx import Presentation

def concatenate_text_runs(paragraph):
    text = ""
    for run in paragraph.runs:
        text += run.text
    return text

def main():
    presentation = Presentation(rf'F:\pptx2md\创意信件毕业设计答辩PPT模板.pptx')
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    text = concatenate_text_runs(paragraph)
                    # You can do something with the concatenated text here
                    # For example, print it or store it in a list or file
                    print(text)

if __name__ == "__main__":
    main()
