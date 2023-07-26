from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

def get_textbox_text_and_position(slide):
    textbox_data = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
            # text = shape.text
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
        else:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                # text = shape.text
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

    return textbox_data


def find_textbox_by_position(slide, left, top):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if shape.left == left and shape.top == top:
                return shape
        else:
            if shape.left == left and shape.top == top:
                return shape

    return None

def replace_text_in_textbox(slide, left, top, target_text, replacement_text):
    textbox = find_textbox_by_position(slide, left, top)
    if textbox:
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                if target_text in run.text:
                    run.text = run.text.replace(target_text, replacement_text)
                    run.font.size = Pt(14)  # 设置字体
                    run.font.bold = False    # 设置字体

# 打开PPT文件
ppt = Presentation(rf'F:\pythonProject7\PPT_File\test000.pptx')

# # 获取第一个幻灯片（可以根据需要选择特定的幻灯片）
# slide = ppt.slides[1-1]
def read_replacement_text_from_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:

        return file.readlines()


replacement_file_path = rf"F:\pptx2md\pptx2md\test\test000_content.txt"
replacement_list = read_replacement_text_from_file(replacement_file_path)
#print(replacement_list)
num = 0
for i, slide in enumerate(ppt.slides, start=1):
    # 获取文本框信息
    textbox_data = get_textbox_text_and_position(slide)
    # 输出文本框信息
    for data in textbox_data:
        if len(data['text']) <= 2:
            continue
        print(f"Left: {data['left']} | Top: {data['top']} | Width: {data['width']} | Height: {data['height']} | Text: {data['text']} | ReplaceText: {replacement_list[num]}")
        # 设置要替换的位置信息
        target_left = int(data['left'])
        target_top =  int(data['top'])
        # 替换文本框中的内容
        target_text = data['text'].strip()
        replace_text_in_textbox(slide, target_left, target_top, target_text,replacement_list[num].strip())
        num+=1


# 保存修改后的PPT文件
ppt.save('new_.pptx')
