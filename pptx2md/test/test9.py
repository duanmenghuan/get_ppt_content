from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_textbox_placeholders_info(presentation_file):
    presentation = Presentation(presentation_file)
    placeholders_info = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                placeholder = shape.text_frame
                index = placeholder.placeholder_format.idx  # 索引
                name = placeholder.placeholder_format.ph_idx  # 名称
                type_ = shape.text_frame.text  # 类型
                placeholders_info.append((index, name, type_))

    return placeholders_info

if __name__ == "__main__":
    presentation_file = rf'F:\pptx2md\test010.pptx'
    placeholders_info = get_textbox_placeholders_info(presentation_file)
    for index, name, type_ in placeholders_info:
        print(f"Index: {index}, Name: {name}, Type: {type_}")
