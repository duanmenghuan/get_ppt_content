from pptx import Presentation
presentation = Presentation(rf'F:\pptx2md\test010.pptx')
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.shape_type == 17:  # 17 表示文本框
            print("索引:", shape.index)
            print("名称:", shape.name)
            print("类型:", shape.shape_type)

