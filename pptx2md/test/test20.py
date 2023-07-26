from pptx import Presentation


def get_font_sizes_with_text(pptx_file):
    prs = Presentation(pptx_file)
    font_sizes_with_text = {}
    slide_num = 0

    for slide in prs.slides:
        slide_num += 1  # Increment the slide number for each new slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            font_size = run.font.size.pt
                            text = run.text.strip()
                            if text:
                                font_sizes_with_text[text] = (font_size, slide_num)  # Include the slide number

    return font_sizes_with_text


if __name__ == "__main__":
    pptx_file_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'
    font_sizes_with_text = get_font_sizes_with_text(pptx_file_path)
    print("Font Sizes, Text, and Slide Number used in the PowerPoint presentation:")
    for text, (size, slide_num) in font_sizes_with_text.items():
        print(f"Slide {slide_num} - Text: '{text}' - Font Size: {size} pt")
