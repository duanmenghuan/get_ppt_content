# from pptx import Presentation
#
# def get_font_size(p):
#     # 获取段落的所有runs
#     runs = p.runs
#     #if len(runs) >= 2:
#         # 获取前两个runs的字体大小
#     font_size1 = runs.font.size
#     #     font_size2 = runs[1].font.size
#     #     # 确保前两个字体大小相同
#     #     if font_size1 == font_size2:
#     #         return font_size1
#     # # 默认情况下，返回整个段落的字体大小
#     # return p.font.size
#     return font_size1
#
# # 读取PPT
# ppt_file_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'
# presentation = Presentation(ppt_file_path)
#
# # 遍历每个幻灯片
# for slide_number, slide in enumerate(presentation.slides, start=1):
#     # 遍历每个形状（包含文本）在幻灯片上
#     for shape in slide.shapes:
#         if shape.has_text_frame:
#             text_frame = shape.text_frame
#             # 获取第一个段落
#             paragraph = text_frame.paragraphs[0]
#             # 获取前两个字的字体大小
#             font_size = get_font_size(paragraph)
#             # 打印输出段落文字和字体大小以及所在的幻灯片编号
#             print(f"Slide {slide_number}: {paragraph.text} (Font size: {font_size})")


from pptx import Presentation
from pptx.util import Pt



def get_font_size(p):
    #获取段落的所有runs
    runs = p.runs
    font_sizes = set(run.font.size for run in runs )
    # 如果段落中所有的run都具有相同的字体大小，则返回该字体大小；否则返回None
    if len(font_sizes) == 1:
        return font_sizes.pop()
    else:
        return None


# 读取PPT
ppt_file_path = rf'F:\pptx2md\test002.pptx'
presentation = Presentation(ppt_file_path)

# 遍历每个幻灯片
for slide_number, slide in enumerate(presentation.slides, start=1):
    # 遍历每个形状（包含文本）在幻灯片上
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            # 获取第一个段落
            paragraph = text_frame.paragraphs[0]
            # 获取前两个字的字体大小
            font_size = get_font_size(paragraph)
            # 打印输出段落文字和字体大小以及所在的幻灯片编号
            if font_size == None:
                continue
            print(f"Slide {slide_number}: {paragraph.text} (Font size: {font_size.pt})")
