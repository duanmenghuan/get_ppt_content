from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

def get_textbox_text_and_position(slide):
    textbox_data = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
            # text = shape.text
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
        else:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                # text = shape.text
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})



    return textbox_data


def find_textbox_by_position(slide, left, top):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if shape.left == left and shape.top == top:
                return shape
    return None

def replace_text_in_textbox(slide, left, top, target_text, replacement_text):
    textbox = find_textbox_by_position(slide, left, top)
    if textbox:
        for paragraph in textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                if target_text in run.text:
                    run.text = run.text.replace(target_text, replacement_text)
                    # run.font.size = Pt(14)  # 设置字体
                    # run.font.bold = False    # 设置字体


# 打开PPT文件
ppt = Presentation(rf'F:\pptx2md\get_ppt\艺术.pptx')
# # 获取第一个幻灯片（可以根据需要选择特定的幻灯片）
slide = ppt.slides[4-1]

    # 获取文本框信息
textbox_data = get_textbox_text_and_position(slide)
# 输出文本框信息
for data in textbox_data:
    if len(data['text']) <= 2:
        #print(data['text'])
        continue
    print(f"Left: {data['left']} | Top: {data['top']} | Width: {data['width']} | Height: {data['height']} | Text: {data['text']}")

#  # 设置要替换的位置信息
# target_left = 1259111
# target_top =  182798
# # 替换文本框中的内容
# target_text = '标题一下的分标题'
# replace_text_in_textbox(slide, target_left, target_top, target_text,'智能行为和思维的方法，')
# # 保存修改后的PPT文件
# ppt.save('new_.pptx')
#
#
# data = [
#     {'Left': 4956466, 'Top': 799966, 'Width': 2199661, 'Height': 1349635, 'Text': '单击填加问题三？'},
#     {'Left': 1158701, 'Top': 2359584, 'Width': 2503679, 'Height': 1360524, 'Text': '单击填加问题一？'},
#     {'Left': 5915713, 'Top': 2226241, 'Width': 2329588, 'Height': 1512917, 'Text': '单击填加问题四？'},
#     {'Left': 2149301, 'Top': 844784, 'Width': 2639616, 'Height': 1359420, 'Text': '单击填加问题二？'},
#     {'Left': 641111, 'Top': 316677, 'Width': 3295650, 'Height': 315461, 'Text': '在此输入您的标题内容'}
# ]
#
# left_values = [entry['Left'] for entry in data]
# minimum_left = min(left_values)
#
# print("The minimum 'Left' value is:", '首页标题')
