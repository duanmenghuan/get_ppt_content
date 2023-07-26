import re
from pptx import Presentation
from pptx.util import Inches

def replace_text_in_ppt(ppt_file, replacements):
    presentation = Presentation(ppt_file)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 替换文本
                        for original_text, replace_text in replacements.items():
                            original_text = original_text.strip()
                            original_text_pattern = re.escape(original_text)
                            print(original_text_pattern)
                            print(replace_text)
                            run.text = re.sub(original_text_pattern, replace_text, run.text)

    presentation.save('modified_presentation1.pptx')

if __name__ == "__main__":
    # 读取输入和替换内容到字典
    replacements = {}
    with open('input1.txt', encoding='UTF8') as input_md:
        with open('replace1.txt', encoding='UTF8') as replace_md:
            for input_line, replace_line in zip(input_md, replace_md):
                original_text = input_line.strip()
                replace_text = replace_line.strip()
                replacements[original_text] = replace_text

    # 调用替换函数
    replace_text_in_ppt(fr'F:\pptx2md\test007.pptx', replacements)
