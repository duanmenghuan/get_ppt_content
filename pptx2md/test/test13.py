from pptx import Presentation

def set_slide_title_and_subtitle(prs, slide_index, title_text, subtitle_text):
    slide = prs.slides[slide_index]
    placeholders = slide.shapes.placeholders

    # 设置标题
    title_placeholder = placeholders[0]
    title_placeholder.text = title_text

    # 设置副标题
    subtitle_placeholder = placeholders[1]
    subtitle_placeholder.text = subtitle_text

# 创建一个新的 PowerPoint 文档
presentation = Presentation()

# 添加新幻灯片
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)

# 设置幻灯片标题和副标题
set_slide_title_and_subtitle(presentation, 0, "这是标题", "这是副标题")

# 保存 PowerPoint 文档
presentation.save("example.pptx")
