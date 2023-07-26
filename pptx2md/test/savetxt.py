#!/usr/bin/env python
# -*- coding:utf-8 -*-
# @Time  : 2023/7/19 9:32
# @Author: Jerry
# @File  : savetxt.py.py

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
import os
def extract_text(pres):
    # 创建一个空列表来存储全部文本及其格式信息
    text_info = []
    processed_text_boxes = set()  # 用于存储已处理的文本框
    # 遍历每个幻灯片
    for slide_num, slide in enumerate(pres.slides, start=1):
        #遍历每个形状
        for shape_num, shape in enumerate(slide.shapes, start=1):
            #shape文本部分 to do pic,title,table,excel
            if shape.has_text_frame:
                # 获取文本框
                text_frame = shape.text_frame
                if text_frame in processed_text_boxes:
                    continue
                # 遍历每个段落
                for paragraph_num, paragraph in enumerate(text_frame.paragraphs, start=1):
                    # 获取段落文本
                    text = paragraph.text.strip()

                    if not text:
                        continue
                    # 为了保存不包含run运行元素的文本框信息，而取消run_num部分
                    # for run_num, run in enumerate(paragraph.runs, start=1):
                    #     if run.font.size is not None:
                    #         # 获取运行元素的字号、位置和颜色等格式信息
                    #         font = run.font
                    #         font_size = font.size.pt
                    #         left = shape.left + text_frame.margin_left
                    #         top = shape.top + text_frame.margin_top
                    #         color = font.color
                    #         if color.type == MSO_COLOR_TYPE.RGB:
                    #             color_type = 'RGB'
                    #             color_value = color.rgb
                    #         elif color.type == MSO_COLOR_TYPE.SCHEME:
                    #             color_type = 'SCHEME'
                    #             color_value = color._color.theme_color
                    #             #print(slide.slide_layout.slide_master.slide_layouts[0])
                    #             #print(dir(pres.slides[0].slide_layout.slide_master.slide_layouts[0].slide_master.shapes))
                    #             #print(RGBColor.from_string(pres.slides[0].slide_layout.slide_master.slide_layouts[0].color_scheme.get_color(MSO_THEME_COLOR.ACCENT_1).rgb))
                    #             #print(MSO_THEME_COLOR.ACCENT_1)
                    #         # 将文本及其格式信息添加到列表中
                    #         text_info.append({
                    #             'location':(slide_num, shape_num, paragraph_num, run_num),
                    #             'text': text,
                    #             'font_size': font_size,
                    #             'left': left,
                    #             'top': top,
                    #             'color_type': color_type,
                    #             'color_value': color_value
                    #         })
                    # 获取运行元素的字号、位置和颜色等格式信息
                    if hasattr(paragraph.runs[0].font.size, 'pt'):
                        font_size = paragraph.runs[0].font.size.pt
                    else:
                        font_size = None
                    left = shape.left + text_frame.margin_left
                    top = shape.top + text_frame.margin_top

                    # 获取颜色信息
                    if paragraph.runs:
                        color = paragraph.runs[0].font.color
                        if color.type == MSO_COLOR_TYPE.RGB:
                            color_type = 'RGB'
                            color_value = color.rgb
                            #alpha = color.rgb.a

                        elif color.type == MSO_COLOR_TYPE.SCHEME:
                            color_type = 'SCHEME'
                            color_value = color.theme_color
                            #alpha = pres.slide_master.color_map.overrides[0].rgb.a
                            #目前SCHEME的文本只能读到预设的颜色类型，暂时不能区分同等类型不同透明度

                    else:
                        color_type = None
                        color_value = None

                    # 将文本及其格式信息添加到列表中
                    text_info.append({
                        'location': (slide_num, shape_num, paragraph_num),
                        'text': text,
                        'font_size': font_size,
                        'left': left,
                        'top': top,
                        'color_type': color_type,
                        'color_value': color_value,
                        #'alpha': alpha
                    })
    return text_info

# 读取PPT文件
currpwd = os.path.dirname(os.path.abspath(__file__))
filepath = f'{currpwd}/fixtures/test2.pptx'
print(f'filepath:{filepath}')
pres = Presentation(filepath)

# 提取文本及其格式信息
text_info = extract_text(pres)

#placeholders

# 输出文本及其格式信息
for info in text_info:
    print(f"Location: {info['location']}")
    print(f"文本: {info['text']}")
    print(f"字号: {info['font_size']}")
    print(f"位置: 左边距{info['left']}, 上边距{info['top']}")
    print(f"颜色类别: {info['color_type']}\n")
    print(f"颜色信息: {info['color_value']}\n")
#for info in text_info:
#    print(text_info)