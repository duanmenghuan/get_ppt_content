from pptx import Presentation

def extract_text_from_shapes(pptx_file):
    prs = Presentation(pptx_file)

    # 遍历每个幻灯片
    for slide in prs.slides:
        # 遍历每个形状
        for shape in slide.shapes:
            # 检查形状是否包含文本
            if hasattr(shape, "text"):
                text = shape.text
                if text:
                    print(text)



if __name__ == "__main__":
    pptx_file = rf'F:\pptx2md\test002.pptx'
    extract_text_from_shapes(pptx_file)
