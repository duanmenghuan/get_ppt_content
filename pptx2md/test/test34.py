from pptx import Presentation
import copy

def count_chars_in_shape(shape):
    """
    计算形状中可以容纳的汉字数量
    这里假设每个汉字占两个英文字符的位置
    """
    text = shape.text_frame.text
    char_count = len(text)
    return char_count // 2

def max_chars_in_shape(shape):
    """
    使用二分查找法计算形状中最大可以容纳的汉字数量
    """
    test_shape = copy.deepcopy(shape)

    # 检查文本框是否已经自适应调整大小
    if shape.text_frame.auto_size == True:
        return count_chars_in_shape(test_shape)

    max_chars = 0
    left = 0
    right = 1_000  # 设置一个较大的初始值，根据实际情况调整
    while left <= right:
        mid = (left + right) // 2
        test_shape.text_frame.text = "汉" * mid  # 假设每个汉字占两个英文字符的位置
        char_count = count_chars_in_shape(test_shape)
        if char_count > mid:
            # 文本框容纳的汉字数量超过了mid，缩小范围继续二分查找
            left = mid + 1
        else:
            # 文本框容纳的汉字数量未达到mid，缩小范围继续二分查找
            right = mid - 1
            max_chars = mid
    return max_chars

def main(ppt_file_path):
    # 打开PPT文件
    presentation = Presentation(ppt_file_path)

    # 遍历每个幻灯片
    for slide in presentation.slides:
        # 遍历幻灯片中的每个形状
        for shape in slide.shapes:
            # 检查形状是否为文本框
            if shape.has_text_frame:
                if len(shape.text_frame.text) > 2:
                    max_char_count = max_chars_in_shape(shape)
                    print(f"形状类型: {shape.shape_type}, 最大可容纳汉字数量: {max_char_count},'文字‘：{shape.text_frame.text}’")

if __name__ == "__main__":
    ppt_file_path = "F:\pptx2md\大气工作总结计划汇报PPT模板.pptx"  # 替换为你的PPT文件路径
    main(ppt_file_path)
