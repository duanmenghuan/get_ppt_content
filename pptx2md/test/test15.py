from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_textbox_text_and_position(slide):
    textbox_data = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            text = shape.text
            textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

    return textbox_data





def replace_text_in_textbox(slide, target_text, replacement_text):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            shape.text = shape.text.replace(target_text, replacement_text)



# 打开PPT文件
ppt = Presentation(rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx')

# 获取第一个幻灯片（可以根据需要选择特定的幻灯片）
slide = ppt.slides[4]

# 获取文本框信息
textbox_data = get_textbox_text_and_position(slide)

# 输出文本框信息
for data in textbox_data:
    print(f"Left: {data['left']} | Top: {data['top']} | Width: {data['width']} | Height: {data['height']} | Text: {data['text']}")

# 替换文本框中的内容
target_text = "这里可以添加主要内容这里可以添加主要内容这里可以添加主要内容这里可以添加主要内容"
replacement_text = "替换后的文本"
replace_text_in_textbox(slide, target_text, replacement_text)

# 保存修改后的PPT文件
ppt.save('modified_presentation.pptx')
