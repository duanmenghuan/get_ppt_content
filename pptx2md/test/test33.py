import cairo
from pptx import Presentation
from pptx.util import Pt


def calculate_characters_capacity_with_font_metrics(text, width_pt, height_pt, font_size_pts, character_width_pts,
                                                    line_spacing_pts):
    # 创建图像表面以绘制文本和测量字体字形
    surface = cairo.ImageSurface(cairo.FORMAT_ARGB32, int(width_pt), int(height_pt))
    ctx = cairo.Context(surface)

    # 设置字体
    ctx.select_font_face("Arial", cairo.FONT_SLANT_NORMAL, cairo.FONT_WEIGHT_NORMAL)
    ctx.set_font_size(font_size_pts)

    # 获取字体字形
    _, _, _, _, ascent, descent, _ = ctx.font_extents()
    line_height_pts = ascent + descent + line_spacing_pts

    # 计算可以容纳的汉字数量
    max_lines = int(height_pt / line_height_pts)
    max_characters_per_line = int(width_pt / character_width_pts)
    total_max_characters = max_lines * max_characters_per_line

    # 将文本分成行并计算字符数量
    lines = text.split('\n')
    total_characters = 0
    for line in lines:
        total_characters += len(line)
        if total_characters >= total_max_characters:
            break

    return total_characters


def pptx_text_info_with_capacity(pptx_file, character_width_pts, line_spacing_pts):
    prs = Presentation(pptx_file)
    result = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            text = ''

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text.strip() + ' '

            width_pt = text_frame.width
            height_pt = text_frame.height
            font_size_pts = text_frame.text_range.font.size.pt

            # 计算每个文本框中可以容纳的汉字数量
            capacity = calculate_characters_capacity_with_font_metrics(
                text, width_pt, height_pt, font_size_pts, character_width_pts, line_spacing_pts
            )

            result.append({
                'text': text.strip(),
                'width_pt': width_pt,
                'height_pt': height_pt,
                'font_size_pts': font_size_pts,
                'capacity': capacity
            })

    return result


# 示例用法
pptx_file_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'
character_width_pts = 12  # 汉字宽度，假设为12磅
line_spacing_pts = 1.5 * character_width_pts  # 行间距，假设为字体大小的1.5倍

text_info_with_capacity = pptx_text_info_with_capacity(pptx_file_path, character_width_pts, line_spacing_pts)

for info in text_info_with_capacity:
    print(f"Text: '{info['text']}'")
    print(f"Width (pt): {info['width_pt']:.2f}")
    print(f"Height (pt): {info['height_pt']:.2f}")
    print(f"Font size (pts): {info['font_size_pts']:.2f}")
    print(f"Max characters: {info['capacity']}")
    print()
