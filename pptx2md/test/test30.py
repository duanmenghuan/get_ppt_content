from pptx import Presentation
from pptx.util import Pt
import math

def emu_to_pixels(emu):
    # 将EMU转换为像素，假设PPTX使用96dpi分辨率
    return emu / 914400 * 96

def get_font_height(font_size):
    # 计算字体的高度（单位：磅）
    return Pt(font_size).emu / 914400

def calculate_max_chars(pptx_path, font_size, textbox_width):
    prs = Presentation(pptx_path)
    slide = prs.slides[0]  # 假设你想在第一张幻灯片中进行计算

    # 遍历所有形状找到包含文本框的形状
    text_box_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
    if not text_box_shapes:
        raise ValueError("在第一张幻灯片中找不到文本框。")

    # 假设我们只处理第一个文本框
    text_box = text_box_shapes[0].text_frame

    # 获取文本框的高度（单位：磅）
    text_box_height = Pt(text_box.margin_bottom).emu / 914400

    # 获取字体的高度（单位：磅）
    font_height = get_font_height(font_size)

    # 计算可以容纳的最大行数
    max_lines = math.floor(text_box_height / font_height)

    # 获取文本框的宽度（单位：像素）
    textbox_width_pixels = emu_to_pixels(textbox_width)

    # 获取字体的宽度（单位：像素，这里假设所有字符的宽度相同）
    # 字体宽度根据字体类型和大小可以有所不同，这里简单地假设为字号的一半
    font_width = font_size / 2

    # 计算每行可以容纳的字符数
    chars_per_line = math.floor(textbox_width_pixels / font_width)

    # 计算可以容纳的最大字符数
    max_chars = max_lines * chars_per_line

    return max_chars

if __name__ == "__main__":
    pptx_file_path = rf"F:\pptx2md\大气工作总结计划汇报PPT模板.pptx"  # 替换为你的PPTX文件路径
    font_size = 12  # 替换为你的字体大小
    textbox_width = 2471314  # 替换为你的文本框宽度（假设为400磅）

    max_chars = calculate_max_chars(pptx_file_path, font_size, textbox_width)
    print(f"文本框（宽度：{textbox_width}px）可容纳约 {max_chars} 个汉字（字号：{font_size}）")
