from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont

import matplotlib.font_manager

fonts = {f.name for f in matplotlib.font_manager.fontManager.ttflist}
print(fonts)


# 函数：将EMU转换为像素
def emu_to_pixels(emu, dpi):
    return emu / 914400 * dpi

# 函数：获取汉字的平均高度（根据实际情况调整）
def average_chinese_character_height(font_size):
    # 这里假设汉字高度约为字号的1.2倍
    return font_size * 1.2

# 函数：获取汉字的平均宽度（根据实际情况调整）
def average_chinese_character_width(font_size):
    # 这里假设汉字宽度约为字号的0.6倍
    return font_size * 0.6

# 函数：渲染文本并获取图像的尺寸
def get_text_dimensions(text, font_size):
    font = ImageFont.truetype("SimHei.ttf", font_size)  # 替换为你的字体文件路径
    image = Image.new("RGB", (1, 1), color="white")
    draw = ImageDraw.Draw(image)
    text_width, text_height = draw.textsize(text, font)
    return text_width, text_height

# 读取PPTX文件
def get_pptx_text_dimensions(file_path):
    presentation = Presentation(file_path)
    total_height_emu = 0
    max_width_emu = 0

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                font_size = 14  # 字体大小，可以根据实际情况调整
                text_width, text_height = get_text_dimensions(text, font_size)

                total_height_emu += Inches(text_height).emu
                width_emu = Inches(text_width).emu
                if width_emu > max_width_emu:
                    max_width_emu = width_emu

    return total_height_emu, max_width_emu

def main():
    file_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'  # 替换为你的PPTX文件路径
    dpi = 96  # 屏幕DPI，可以根据实际情况调整

    text_height_emu, max_width_emu = get_pptx_text_dimensions(file_path)
    text_height_pixels = emu_to_pixels(text_height_emu, dpi)
    max_width_pixels = emu_to_pixels(max_width_emu, dpi)

    font_size = 14  # 字体大小，可以根据实际情况调整
    chinese_character_height = average_chinese_character_height(font_size)
    chinese_character_width = average_chinese_character_width(font_size)

    # 计算可以容纳的汉字数量
    num_of_chinese_characters_height = int(text_height_pixels / chinese_character_height)
    num_of_chinese_characters_width = int(max_width_pixels / chinese_character_width)
    num_of_chinese_characters = num_of_chinese_characters_height * num_of_chinese_characters_width

    print("PPTX文本高度（像素）：", text_height_pixels)
    print("PPTX文本最大宽度（像素）：", max_width_pixels)
    print("平均每个汉字的高度（像素）：", chinese_character_height)
    print("平均每个汉字的宽度（像素）：", chinese_character_width)
    print("可以容纳的汉字数量（仅考虑高度和宽度）：", num_of_chinese_characters)

if __name__ == "__main__":
    main()
