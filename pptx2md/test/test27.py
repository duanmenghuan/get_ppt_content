from pptx import Presentation

ppt = Presentation("1.pptx")
for slide in ppt.slides:			#> .slides 得到一个列表，包含每个列表slide
	#print(slide)
	for shape in slide.shapes:		#> slide.shapes 形状
		if shape.has_text_frame:	#shape.has_text_frame 判断是否有文字
			text_frame = shape.has_text_frame #shape.text_frame 获取文字框
			#print(text_frame.text)
			for paragraph in text_frame.paragraphs:	#text_frame.paragraphs 获取段落
				print(paragraph.text)
