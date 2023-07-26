from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_boxes(pptx_file):
    prs = Presentation(pptx_file)
    text_boxes = []

    for i, slide in enumerate(prs.slides, start=1):
        num = 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                if text:
                    text = str(num)+':'+text
                    text_boxes.append((i, text))
                    num+=1
                #if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    print(left)
                    print(top)
                    print(width)
                    print(height)


    return text_boxes

pptx_file = rf'F:\pptx2md\test007.pptx'

text_boxes = extract_text_boxes(pptx_file)

for slide_num, text_box in text_boxes:
    print(f"Slide {slide_num}: {text_box}")

