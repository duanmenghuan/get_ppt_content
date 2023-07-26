from pptx import Presentation

#函数名为extract_text_boxes，它有一个参数pptx_file，该参数用于传入PPTX文件的路径。
def extract_text_boxes(pptx_file):
    i = 0
    #这一行创建了一个Presentation对象，用于打开PPTX文件并加载其内容。Presentation是python-pptx库中的一个类，它允许我们读取和操作PowerPoint文档。
    prs = Presentation(pptx_file)
    #创建一个空列表text_boxes，用于存储提取的文本内容
    text_boxes = []
    #这是一个外部循环，遍历Presentation对象中的每一张幻灯片
    for slide in prs.slides:
        i+=1
        #for shape in slide.shapes:这是一个内部循环，遍历每张幻灯片中的每个形状（shape）
        for shape in slide.shapes:
            #这个条件判断语句用于检查形状是否包含文本框
            if shape.has_text_frame:
                #这一行获取文本框对象，以便后续从中提取文本。
                text_frame = shape.text_frame
                #这一行使用列表推导式（list comprehension）提取文本框中每个段落的文本，并去除文本两端的空格。
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                if text:
                    text_boxes.append(text)


    return text_boxes

# 指定PPTX文件路径
pptx_file = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'

# 提取文本框内容
text_boxes = extract_text_boxes(pptx_file)

# 打印文本框内容
for text_box in text_boxes:
    print(text_box)

'''
先了解下PPT基本结构在python分别是什么含义：
Slide：幻灯片，就是演示文稿中每一页的页面。
Shape：方框，在每页幻灯片内插入的方框，可以是形状，也可以是文本框。
Run：文字块，一般为较少字符。
Paragraph：段落，通常有序号ㆍ、1.等。
'''
