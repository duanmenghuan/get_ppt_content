from pptx import  Presentation
prs = Presentation(rf'F:\pptx2md\test010.pptx')
shape = prs.slides.add_slide(prs.slide_layouts[0])
# for shape in prs.slides:
#     for placeholder in shape.placeholders:
#         message = placeholder.placeholder_format
#         print(f'索引:{message.idx},名称：{placeholder.name},类型:{message.type}')


title = shape.placeholders[0]
subtitle = shape.placeholders[1]
title.text = '通过标题索引修改'
subtitle.text = '通过副标题索引修改'
prs.save('test.pptx')

# from pptx import Presentation
#
# prs = Presentation(rf'F:\pptx2md\test010.pptx')
#
# for slide in prs.slides:
#     for shape in slide.shapes:
#         if shape.is_placeholder:
#             placeholder = shape.placeholder_format
#             print(f'索引: {placeholder.idx}, 名称: {shape.name}, 类型: {placeholder.type}')

