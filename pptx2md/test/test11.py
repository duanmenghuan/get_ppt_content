from pptx import Presentation

def get_text_info(presentation_file):
    prs = Presentation(presentation_file)
    text_info = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for  paragraph in shape.text_frame.paragraphs:
                 # for idx, run in enumerate(paragraph.runs):
                    if paragraph.text == " ":
                        continue
                    text = paragraph.text
                    name = shape.name
                    shape_type = shape.shape_type
                    text_info.append({
                        '幻灯片页数': slide_number,
                        '形状': name,
                        '形状类型': shape_type,
                        '内容': text,
                        # '索引': idx,
                    })

    return text_info

if __name__ == "__main__":
    pptx_file = rf"F:\pptx2md\大气工作总结计划汇报PPT模板.pptx"
    text_info_list = get_text_info(pptx_file)

    for info in text_info_list:
        print(info)
