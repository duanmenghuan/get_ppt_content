from pptx import Presentation
prs = Presentation(rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx') # 打开一个已有的 PPT 文件

for slide in prs.slides:
    for shape in  slide.shapes:
        if shape.has_text_frame:  # 判断是否有文本框
            text_frame = shape.text_frame  # 获取文本框中的文本帧对象
            #print(text_frame)
            #paragraph = text_frame.paragraphs  # 获取第一个段落对象
            for paragraph in text_frame.paragraphs:
                level = paragraph.level  # 获取段落级别
                if level == 0:
                    if len(paragraph.text.strip()) <= 1:
                        continue
                    if '标题' in paragraph.text.strip() and len(paragraph.text.strip())<20 : #
                        print(f"{paragraph.text.strip()}","此内容为标题")
                    elif '文字'in paragraph.text.strip() or '文本' in paragraph.text.strip() or '内容' in paragraph.text.strip() :
                        print(f"{paragraph.text.strip()}", "此内容为正文")
                    else:
                        print(f"{paragraph.text.strip()}")
                elif level == 1:
                    print(f"{paragraph.text.strip()}","副标题")
                else:
                    print(f"{paragraph.text.strip()}","正文")



