from pptx import Presentation
from pptx.util import Inches

presentation = Presentation(fr'F:\pptx2md\test010.pptx')

# 读取输入和替换内容到列表
with open(f'input.txt', encoding='UTF8') as input_md:
    input_lines = input_md.readlines()

with open(f'replace.txt', encoding='UTF8') as replace_md:
    replace_lines = replace_md.readlines()


print(len(input_lines))
print(len(replace_lines))
print(input_lines)
print(replace_lines)
# 遍历每个幻灯片并替换内容
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            # for paragraph in shape.text_frame.paragraphs:
            #     for run in paragraph.runs:
            if shape.has_text_frame:
                    #run.text = run.text.replace('标题一下的正文标题一下的正文内容一', '计算机视觉是一门研究如何让计算机理解和处理图像和视频的科学。\n')
                    # # 替换文本
                    for i in range(len(input_lines)):
                        #if input_lines[i].strip() in run.text:
                        print(input_lines[i])
                        print( replace_lines[i])
                        shape.text = shape.text.replace(input_lines[i].strip(), replace_lines[i].strip())

presentation.save('modified_presentation.pptx')


