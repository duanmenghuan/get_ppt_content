from pptx import Presentation

def extract_text_boxes(pptx_file):
    prs = Presentation(pptx_file)
    text_boxes = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                #这是一个循环语句，遍历文本框中的每个段落。在PowerPoint中，文本框的内容可以被分成多个段落，每个段落可能包含不同的格式设置（例如字体、颜色等）。
                for paragraph in text_frame.paragraphs:
                    #这是一个内部循环，遍历当前段落中的每个运行（run）。运行是一系列具有相同格式设置的连续文本。
                    for run in paragraph.runs:
                        #这一行从当前运行中获取文本内容，并使用strip()方法去除文本两端的空格。strip()方法是字符串方法，它会删除字符串开头和结尾的空格和换行符
                        text = run.text.strip()
                        if text:
                            text_boxes.append(text)
    text_boxes.append("\t")

    return text_boxes

# 指定PPTX文件路径
pptx_file = 'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'

# 提取文本框内容
text_boxes = extract_text_boxes(pptx_file)

# 打印文本框内容
for text_box in text_boxes:
    print(text_box)

'''
先了解下PPT基本结构在python分别是什么含义：
Slide：幻灯片，就是演示文稿中每一页的页面。
Shape：方框，在每页幻灯片内插入的方框，可以是形状，也可以是文本框。
Run：文字块，一般为较少字符。
Paragraph：段落，通常有序号ㆍ、1.等。
'''

