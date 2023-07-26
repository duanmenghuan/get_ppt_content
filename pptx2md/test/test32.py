from pptx import Presentation
from pptx.util import Pt


def pptx_text_info(pptx_file):
    prs = Presentation(pptx_file)
    result = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            text = ''
            font_size_pts = None

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text.strip() + ' '
                    # Check if the font size is not None before accessing 'pt'
                    if run.font.size is not None:
                        font_size_pts = run.font.size.pt if font_size_pts is None else max(font_size_pts,
                                                                                           run.font.size.pt)

            width_pt = shape.width
            height_pt = shape.height

            result.append({
                'text': text.strip(),
                'width_pt': width_pt,
                'height_pt': height_pt,
                'font_size_pts': font_size_pts or 12  # Use a default font size of 12 if font_size_pts is None
            })

    return result


def calculate_characters_capacity(text_info, character_width_pts):
    # 计算可以容纳多少汉字
    for info in text_info:
        max_characters = int(info['width_pt'] / character_width_pts)
        info['max_characters'] = max_characters

    return text_info


# 示例用法
pptx_file_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'
character_width_pts = 12  # 假设每个汉字宽度为12磅

text_info = pptx_text_info(pptx_file_path)
text_info_with_capacity = calculate_characters_capacity(text_info, character_width_pts)

for info in text_info_with_capacity:
    print(f"Text: '{info['text']}'")
    print(f"Width (pt): {info['width_pt']:.2f}")
    print(f"Height (pt): {info['height_pt']:.2f}")
    print(f"Font size (pts): {info['font_size_pts']:.2f}")
    print(f"Max characters: {info['max_characters']}")
    print()
