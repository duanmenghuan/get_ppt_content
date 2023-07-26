# 导入python-pptx库
from pptx import Presentation
# 导入unicodedata库
import unicodedata

# 打开PPT文件
ppt = Presentation(rf"F:\pptx2md\大气工作总结计划汇报PPT模板.pptx")
# 遍历每一张幻灯片
for slide in ppt.slides:
    # 获取幻灯片的编号
    slide_number = slide.slide_id
    # 打印幻灯片的编号
    print(f"幻灯片{slide_number}：")
    # 遍历幻灯片中的每一个形状
    for shape in slide.shapes:
        # 获取形状的类型
        shape_type = shape.shape_type
        # 打印形状的类型
        print(f"- 形状类型：{shape_type}")
        # 判断是否是文本框对象
        if shape.has_text_frame:
            # 获取文本框对象
            text_frame = shape.text_frame
            # 获取文本内容
            text = text_frame.text

            # 计算形状的面积（单位是英寸）
            area = shape.width * shape.height
            # 计算文本框的有效面积（减去边距）
            effective_area = (shape.width - text_frame.margin_left - text_frame.margin_right) * (shape.height - text_frame.margin_top - text_frame.margin_bottom)
            # 获取文字的字体和大小（假设整个文本框使用同一种字体和大小）
            try: # 尝试获取字体和大小
                font = text_frame.paragraphs[0].runs[0].font
                font_name = font.name
                font_size = font.size
            except IndexError: # 如果出现索引错误
                print("无法获取字体和大小，请检查文本框是否为空或者段落是否有文本")
                continue # 跳过该循环，继续下一个循环

            # 定义一个函数，根据汉字的编码判断其宽度（全角或半角）
            def get_width(char):
                # 如果是汉字，返回2
                if unicodedata.east_asian_width(char) in ["F", "W"]:
                    return 2
                # 否则，返回1
                else:
                    return 1

            # 计算文本的字符数和宽度（单位是磅）
            char_count = 0 # 字符数
            text_width = 0 # 文本宽度
            for char in text:
                # 累加字符数
                char_count += 1
                # 累加字符宽度（根据全角或半角乘以字体大小的一半）
                text_width += get_width(char) * font_size / 2

            # 计算形状里面可以容纳的最大字符数（假设每行只有一个字符，忽略换行符）
            max_char_count = int(effective_area / font_size)
            # 计算形状里面的空余空间（单位是磅）
            free_space = effective_area * 72 - text_width # 72磅等于1英寸

            # 打印结果
            print(f"- 形状的大小是{shape.width:.2f} x {shape.height:.2f} 英寸，面积是{area:.2f} 平方英寸")
            print(f"- 文本框的边距是{text_frame.margin_left:.2f}, {text_frame.margin_right:.2f}, {text_frame.margin_top:.2f}, {text_frame.margin_bottom:.2f} 磅")
            print(f"- 文本框的有效面积是{effective_area:.2f} 平方英寸")
            print(f"- 文字的字体是{font_name}，大小是{font_size} 磅")
            print(f"- 文本的内容是{text}")
            print(f"- 文本的字符数是{char_count}")
            print(f"- 文本的宽度是{text_width:.2f} 磅")
            print(f"- 形状里面可以容纳的最大字符数是{max_char_count}")
            print(f"- 形状里面的空余空间是{free_space:.2f} 磅")
        # 判断是否是GroupShape对象
        elif shape.shape_type == 6: # 6是GroupShape对象的类型编号
            # 获取GroupShape对象内部的形状列表
            subshapes = shape.shapes
            # 遍历形状列表，找到文本框对象
            for subshape in subshapes:
                # 判断是否是文本框对象
                if subshape.has_text_frame:
                    # 获取文本框对象
                    text_frame = subshape.text_frame
                    # 获取文本内容
                    text = text_frame.text
                    # 打印文本内容
                    print(f"- 文本内容：{text}")
