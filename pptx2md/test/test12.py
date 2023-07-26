from pptx import Presentation

# 读取 PowerPoint 文件
presentation = Presentation(rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx')

# 遍历每一张幻灯片
for slide in presentation.slides:
    # 检查是否有占位符
    if slide.placeholders:
        # 如果有占位符，则输出占位符的索引和名称
        for idx, placeholder in enumerate(slide.placeholders):
            print(f"Slide {slide.slide_id}, Placeholder {idx}: {placeholder.name}")
    else:
        print(f"Slide {slide.slide_id} has no placeholders.")
