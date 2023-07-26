from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

def get_slide_content(slide):
    content = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text.strip()
                if text:
                    content.append((text, shape.text_frame.text, shape.name))
    return content

def determine_content_type(slide_content):
    title_keywords = ["title", "heading", "主题"]
    subtitle_keywords = ["subtitle", "副标题"]
    for text, _, shape_name in slide_content:
        if any(keyword in shape_name.lower() for keyword in title_keywords):
            return "标题", text
        elif any(keyword in shape_name.lower() for keyword in subtitle_keywords):
            return "副标题", text
    return "正文", [text for text, _, _ in slide_content]

if __name__ == "__main__":
    ppt_path = rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'

    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        slide_content = get_slide_content(slide)
        content_type, content = determine_content_type(slide_content)
        print(f"第 {i+1} 页 - {content_type}:\n{content}\n")
