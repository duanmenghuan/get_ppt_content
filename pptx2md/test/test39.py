from pptx import Presentation

prs = Presentation('F:\pptx2md\大气工作总结计划汇报PPT模板.pptx') # 打开一个PPT文件
for slide in prs.slides: # 遍历每张幻灯片
    for shape in slide.shapes: # 遍历每个形状
      if hasattr(shape, "text"): # 检查是否有文本框
        text = shape.text # 获取文本框内容
        count = len(text) # 获取文本框字数
        if count == 0: # 如果文本框为空，就跳过这个文本框
          continue
        width = shape.width # 获取文本框宽度
        height = shape.height # 获取文本框高度
        margin_left = shape.text_frame.margin_left # 获取左边距
        margin_right = shape.text_frame.margin_right # 获取右边距
        margin_top = shape.text_frame.margin_top # 获取上边距
        margin_bottom = shape.text_frame.margin_bottom # 获取下边距

        print(width,height,margin_left,margin_right,margin_top,margin_bottom)
        font_size = shape.text_frame.paragraphs[0].font.size # 获取字体大小
        if font_size is None: # 如果字体大小没有值，就给它一个默认值
          font_size = 14
        line_spacing = shape.text_frame.paragraphs[0].line_spacing # 获取行距
        if line_spacing is None: # 如果行距没有值，就给它一个默认值
          line_spacing = 1.5


        area = (width - margin_left - margin_right) * (height - margin_top - margin_bottom) # 计算有效面积
        #print("有效面积",area)
        char_area = font_size ** 2 * 0.3 # 计算每个字符的平均面积（假设中文字符）
        #print("占地面积",char_area*91400)
        chars_per_line = int((width - margin_left - margin_right) // (font_size * 0.3)) # 计算每行的字符数
        lines_per_page = int((height - margin_top - margin_bottom) // (font_size + line_spacing)) # 计算每页的行数
        max_chars = (chars_per_line) / 91400 * (lines_per_page / 91400)*10 # 计算最多的字符数
        print("文字:",text)
        print("文字数量:", count) # 打印文字数量
        print("最多的字符数:", round(max_chars, 0)) # 打印最多的字符数（以字数为单位）
        #print("比例:", count / max_chars) # 打印两者的比例
        print("**************")
