from pptx import Presentation

def get_all_textboxes_content(slide):
    text_boxes_content = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # 获取文本框内容
        if shape.text_frame:
            text_content = shape.text_frame.text
            text_boxes_content.append(text_content)

    return text_boxes_content

def main():
    # 读取pptx文件
    presentation = Presentation(rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx')

    for slide in presentation.slides:
        content = get_all_textboxes_content(slide)
        if content:
            print("Slide #{}:".format(slide.slide_id))
            for idx, text in enumerate(content):
                print("Text Box #{} Content:".format(idx + 1))
                print("Text:", text)
                print()

if __name__ == "__main__":
    main()
