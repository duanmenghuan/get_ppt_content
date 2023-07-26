from pptx import Presentation

def get_text_from_shape(shape):
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""

def get_text_from_pptx(pptx_file_path):
    presentation = Presentation(pptx_file_path)
    text_list = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            text = get_text_from_shape(shape)
            if text:
                text_list.append(text)

    return text_list

if __name__ == "__main__":
    pptx_file_path = rf'F:\pptx2md\test002.pptx'
    text_list = get_text_from_pptx(pptx_file_path)
    for idx, text in enumerate(text_list):
        print(f"Text {idx + 1}: {text}")
