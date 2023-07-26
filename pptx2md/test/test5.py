import re
from pptx import Presentation
from pptx.util import Inches

def replace_text_in_ppt(ppt_file, replacements):
    presentation = Presentation(ppt_file)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 替换文本
                        for original_text, replace_text in replacements.items():
                            if original_text in run.text:
                                run.text = re.sub(rf"\b{re.escape(original_text)}\b", replace_text, run.text)

    presentation.save('modified_presentation.pptx')

if __name__ == "__main__":
    # 读取输入和替换内容到字典
    replacements = {}
    with open('input.txt', encoding='UTF8') as input_md:
        with open('replace.txt', encoding='UTF8') as replace_md:
            for input_line, replace_line in zip(input_md, replace_md):
                original_text = input_line.strip()
                replace_text = replace_line.strip()
                replacements[original_text] = replace_text

    # 调用替换函数
    replace_text_in_ppt(fr'F:\pptx2md\test010.pptx', replacements)
