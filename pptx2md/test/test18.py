from pptx import Presentation
from pptx.dml.color import RGBColor

def get_text_properties(text_run):
    font = text_run.font
    font_color = font.color
    if font_color is None or font_color.type == "not a color":
        color = None
    elif font_color.type == RGBColor:
        color = font_color.rgb
    else:
        color = font_color.theme_color
        if color is None:
            raise Exception("演示文稿中未设置主题颜色。")
    font_size = font.size
    font_size_pt = font_size.pt if font_size is not None else None
    return {
        "text": text_run.text,
        "font_name": font.name,
        "font_size": font_size_pt,
        "color": color
    }

def group_text_runs_by_color(paragraph):
    groups = []
    current_group = []
    current_color = None

    for run in paragraph.runs:
        properties = get_text_properties(run)
        color = properties["color"]

        if color == current_color:
            current_group.append(properties)
        else:
            if current_group:
                groups.append((current_color, current_group))
            current_color = color
            current_group = [properties]

    if current_group:
        groups.append((current_color, current_group))

    return groups

def extract_text_properties(slide):
    text_properties = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            text_runs_by_color = group_text_runs_by_color(paragraph)
            for color, text_runs in text_runs_by_color:
                text_properties.append({"color": color, "text_runs": text_runs})
    return text_properties

def main():
    presentation = Presentation(rf'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx')
    for slide in presentation.slides:
        try:
            slide_properties = extract_text_properties(slide)
            for prop in slide_properties:
                color = prop["color"]
                print(f"Color: {color}")
                for run in prop["text_runs"]:
                    text = run["text"]
                    font_name = run["font_name"]
                    font_size = run["font_size"]
                    print(f"Text: {text}, Font: {font_name}, Font Size: {font_size}")
        except Exception as e:
            print(f"第 {presentation.slides.index(slide) + 1} 页出现异常：{e}")

if __name__ == "__main__":
    main()
