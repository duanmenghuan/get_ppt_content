from pptx import Presentation

def extract_text_boxes(pptx_file):
    prs = Presentation(pptx_file)
    text_boxes = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                if text:
                    text_boxes.append(text)

    return text_boxes

def save_to_md_file(text_boxes, output_file):
    with open(output_file, 'w', encoding='utf-8') as md_file:
        for text_box in text_boxes:
            md_file.write(text_box + '\n')


pptx_file = 'F:\pptx2md\大气工作总结计划汇报PPT模板.pptx'

text_boxes = extract_text_boxes(pptx_file)

output_file = 'output.md'

save_to_md_file(text_boxes, output_file)


