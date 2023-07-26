'''
这段代码包含了一系列导入语句和全局变量定义。
1. `from __future__ import print_function` 是用于将 Python 2 中的 `print` 语句转换为 Python 3 中的 `print` 函数的导入语句。它确保在 Python 2 环境下使用 `print` 函数进行打印。
2. `import collections` 和 `import collections.abc` 导入了 Python 的 `collections` 模块和 `collections.abc` 模块。这些模块提供了各种集合数据类型和抽象基类，用于处理集合对象。
3. `import pptx` 导入了 `pptx` 模块，用于处理 Microsoft PowerPoint 文件。
4. `from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE` 导入了 `pptx` 模块中的枚举类型 `PP_PLACEHOLDER_TYPE` 和 `MSO_SHAPE_TYPE`，用于表示 PowerPoint 中的占位符类型和形状类型。
5. `from pptx.enum.dml import MSO_COLOR_TYPE` 和 `from pptx.enum.dml import MSO_THEME_COLOR_INDEX` 导入了 `pptx` 模块中的枚举类型 `MSO_COLOR_TYPE` 和 `MSO_THEME_COLOR_INDEX`，用于表示 PowerPoint 中的颜色类型和主题颜色索引。
6. `from PIL import Image` 导入了 `PIL`（Python Imaging Library）模块中的 `Image` 类，用于处理图像文件。
7. `import os` 导入了 Python 的 `os` 模块，用于与操作系统进行交互，例如文件路径操作和目录操作。
8. `from rapidfuzz import process as fuze_process` 导入了 `rapidfuzz` 模块中的 `process` 函数，并将其重命名为 `fuze_process`，用于进行模糊匹配和字符串处理。
9. `from operator import attrgetter` 导入了 Python 的 `operator` 模块中的 `attrgetter` 函数，用于获取对象的属性。
10. `from tqdm import tqdm` 导入了 `tqdm` 模块中的 `tqdm` 函数，用于显示进度条。
11. `from pptx2md.global_var import g` 导入了 `pptx2md` 包中的 `global_var` 模块，并从中导入了 `g` 对象。这是一个全局变量对象，用于存储程序运行时的各种设置和状态。
12. `from pptx2md import global_var` 导入了 `pptx2md` 包中的 `global_var` 模块，用于访问其中定义的全局变量。
最后一行的 `global out` 并不是有效的 Python 语法。它可能是一个错误的代码片段或者是被误写在这里的。在正确的 Python 代码中，`global` 关键字用于在函数内部声明一个变量为全局变量，而不是在这里声明。因此，这行代码可能需要删除或者根据实际需要进行修改。
'''
from __future__ import print_function

import collections
import collections.abc
import pptx
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

from PIL import Image
import os
from rapidfuzz import process as fuze_process
from operator import attrgetter

from tqdm import tqdm
from pptx2md.global_var import g
from pptx2md import global_var

picture_count = 0

global out


# pptx type defination rules  定义规则
'''
这段代码定义了一个名为 is_title 的函数，用于判断给定的形状（shape）是否为标题。
以下是代码的解释：
定义了一个函数 is_title，接受一个形状对象 shape 作为输入参数。
判断给定的形状是否是一个占位符（placeholder），并且其占位符类型（placeholder_format.type）为以下之一：
PP_PLACEHOLDER_TYPE.TITLE：标题
PP_PLACEHOLDER_TYPE.SUBTITLE：副标题
PP_PLACEHOLDER_TYPE.VERTICAL_TITLE：垂直标题
PP_PLACEHOLDER_TYPE.CENTER_TITLE：居中标题
如果形状满足上述条件，则返回 True，表示该形状是一个标题。
如果形状不满足上述条件，则返回 False，表示该形状不是一个标题。
该函数用于识别给定形状是否代表一个标题，它检查形状是否是占位符，并且占位符类型与标题相关。
'''
def is_title(shape):
    if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                                 or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE
                                 or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.VERTICAL_TITLE
                                 or shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE):
        return True
    return False

'''
这段代码定义了一个名为 `is_text_block` 的函数，用于判断给定的形状（shape）是否为文本块。
以下是代码的解释：
1. 定义了一个函数 `is_text_block`，接受一个形状对象 `shape` 作为输入参数。
2. 首先判断给定的形状是否具有文本框（text frame），即判断 `shape.has_text_frame` 是否为真。
3. 如果形状是一个占位符，并且占位符类型为 `PP_PLACEHOLDER_TYPE.BODY`，则返回 `True`，表示该形状是一个文本块。
4. 如果形状的文本内容长度大于全局变量 `g.text_block_threshold` 的阈值，则返回 `True`，表示该形状是一个文本块。
5. 如果以上条件都不满足，则返回 `False`，表示该形状不是一个文本块。
该函数用于判断给定形状是否代表一个文本块，它检查形状是否具有文本框，并根据不同的条件判断形状是否属于文本块。
具体条件包括形状是否为占位符类型的主体（body）占位符，以及形状的文本内容长度是否超过了全局变量中设定的阈值。
'''
def is_text_block(shape):
    if shape.has_text_frame:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.BODY:
            return True
        if len(shape.text) > g.text_block_threshold:
            return True
    return False

'''
这段代码定义了一个名为 `is_list_block` 的函数，用于判断给定的形状（shape）是否为列表块。
以下是代码的解释：
1. 定义了一个函数 `is_list_block`，接受一个形状对象 `shape` 作为输入参数。
2. 创建一个空列表 `levels`，用于存储段落的级别。
3. 遍历形状的文本框中的每个段落（`shape.text_frame.paragraphs`）。
4. 如果当前段落的级别（`para.level`）不在 `levels` 列表中，将其添加到 `levels` 列表中。
5. 如果当前段落的级别不为 0 或 `levels` 列表的长度大于 1，则返回 `True`，表示该形状是一个列表块。
6. 如果以上条件都不满足，则返回 `False`，表示该形状不是一个列表块。
该函数用于判断给定形状是否代表一个列表块，它检查形状的文本框中的每个段落的级别，如果存在多个不同的级别或级别不为 0，则判定为列表块。
'''
def is_list_block(shape):
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False

'''
这段代码定义了一个名为 `is_accent` 的函数，用于判断给定的字体（font）是否为强调（accent）样式。
以下是代码的解释：
1. 定义了一个函数 `is_accent`，接受一个字体对象 `font` 作为输入参数。
2. 判断字体是否具有下划线（`font.underline`）或斜体（`font.italic`），如果满足其中一项，则返回 `True`，表示该字体为强调样式。
3. 判断字体颜色是否属于主题颜色索引 `ACCENT_1`、`ACCENT_2`、`ACCENT_3`、`ACCENT_4`、`ACCENT_5` 或 `ACCENT_6`，如果是，则返回 `True`，表示该字体为强调样式。
4. 如果以上条件都不满足，则返回 `False`，表示该字体不是强调样式。
该函数用于判断给定的字体是否属于强调样式，它检查字体的下划线、斜体以及颜色是否符合强调样式的特征。如果满足任一条件，则判定为强调样式。
'''
def is_accent(font):
    if font.underline or font.italic or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and
            (
                    font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_1 or font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_2
                    or font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_3 or font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_4
                    or font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_5 or font.color.theme_color == MSO_THEME_COLOR_INDEX.ACCENT_6)):
        return True
    return False

'''
这段代码定义了一个名为 `is_strong` 的函数，用于判断给定的字体（font）是否为加粗（strong）样式。
以下是代码的解释：
1. 定义了一个函数 `is_strong`，接受一个字体对象 `font` 作为输入参数。
2. 判断字体是否为加粗（`font.bold`），如果是，则返回 `True`，表示该字体为加粗样式。
3. 判断字体颜色是否属于主题颜色索引 `DARK_1` 或 `DARK_2`，如果是，则返回 `True`，表示该字体为加粗样式。
4. 如果以上条件都不满足，则返回 `False`，表示该字体不是加粗样式。
该函数用于判断给定的字体是否属于加粗样式，它检查字体是否为加粗或字体颜色是否属于深色主题颜色索引。如果满足任一条件，则判定为加粗样式。
'''
def is_strong(font):
    if font.bold or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and (font.color.theme_color == MSO_THEME_COLOR_INDEX.DARK_1
                                                          or font.color.theme_color == MSO_THEME_COLOR_INDEX.DARK_2)):
        return True
    return False

'''
这段代码定义了一个名为 `get_formatted_text` 的函数，用于获取格式化的文本内容。
以下是代码的解释：
1. 定义了一个函数 `get_formatted_text`，接受一个段落对象 `para` 作为输入参数。
2. 创建一个空字符串 `res`，用于存储格式化后的文本内容。
3. 遍历段落中的每个文本运行（`para.runs`）。
4. 获取当前文本运行的文本内容（`run.text`）。
5. 如果文本内容为空字符串，则继续下一次循环。
6. 如果全局变量 `g.disable_escaping` 为假（即未禁用转义），则调用 `out.get_escaped` 方法对文本内容进行转义处理。
7. 尝试获取文本运行的超链接地址（`run.hyperlink.address`）。
   - 如果成功获取超链接地址，则调用 `out.get_hyperlink` 方法将文本内容转换为超链接格式。
   - 如果获取超链接地址出现异常（例如解析错误），则使用默认的错误地址调用 `out.get_hyperlink` 方法。
8. 判断文本运行的字体是否属于强调样式（调用 `is_accent(run.font)` 和 `is_strong(run.font)` 函数）。
   - 如果属于强调样式，则调用相应的 `out.get_accent` 或 `out.get_strong` 方法将文本内容转换为强调格式。
9. 如果全局变量 `g.disable_color` 为假（即未禁用颜色），则判断文本运行的字体颜色类型是否为 RGB 类型（`run.font.color.type == MSO_COLOR_TYPE.RGB`）。
   - 如果是 RGB 类型，则调用 `out.get_colored` 方法将文本内容转换为带有颜色的格式。
10. 将格式化后的文本内容追加到 `res` 字符串中。
11. 返回去除首尾空白字符的结果字符串 `res`。
该函数用于将给定段落的文本内容进行格式化处理，根据文本运行的属性（如转义、超链接、强调样式、颜色等）对文本内容进行相应的转换，并返回格式化后的文本结果。
'''
def get_formatted_text(para):
    res = ''
    for run in para.runs:
        text = run.text
        if text == '':
            continue
        if not g.disable_escaping:
            text = out.get_escaped(text)
        try:
            if run.hyperlink.address:
                text = out.get_hyperlink(text, run.hyperlink.address)
        except:
            text = out.get_hyperlink(text, 'error:ppt-link-parsing-issue')
        if is_accent(run.font):
            text = out.get_accent(text)
        elif is_strong(run.font):
            text = out.get_strong(text)
        if not g.disable_color:
            if run.font.color.type == MSO_COLOR_TYPE.RGB:
                text = out.get_colored(text, run.font.color.rgb)
        res += text
    return res.strip()

'''
这段代码定义了一个名为 `process_title` 的函数，用于处理标题形状（shape）。
以下是代码的解释：
1. 定义了一个函数 `process_title`，接受一个标题形状对象 `shape` 和幻灯片索引 `slide_idx` 作为输入参数。
2. 声明 `notes` 列表，用于存储处理过程中的备注信息。
3. 获取标题形状的文本内容，并去除首尾的空白字符（`shape.text_frame.text.strip()`）。
4. 如果全局变量 `g.use_custom_title` 为真，表示使用自定义标题。
   - 使用 `fuze_process.extractOne` 函数从自定义标题列表中找到与当前标题文本最匹配的标题（得分大于等于 96）。
   - 如果找不到匹配的标题，则调用 `out.put_title` 方法将当前标题文本作为新的自定义标题，并传递一个新的标题级别（`g.max_custom_title + 1`）。
   - 如果找到匹配的标题，则将匹配的标题作为新标题，并根据匹配的标题级别调用 `out.put_title` 方法。
     同时，将一条备注信息添加到 `notes` 列表中，表示当前幻灯片的标题已根据自定义标题文件进行转换。
5. 如果全局变量 `g.use_custom_title` 为假，表示不使用自定义标题，直接调用 `out.put_title` 方法将当前标题文本作为一级标题。
6. 返回 `notes` 列表，其中包含了处理过程中的备注信息。
该函数用于处理给定的标题形状，根据是否使用自定义标题来确定如何处理标题文本。如果使用自定义标题，则尝试将当前标题文本与自定义标题列表进行匹配，并根据匹配结果决定新的标题文本和级别。
如果不使用自定义标题，则将当前标题文本作为一级标题处理。最后，返回处理过程中的备注信息列表。
'''
def process_title(shape, slide_idx):
    global out
    notes = []
    text = shape.text_frame.text.strip()
    if g.use_custom_title:
        res = fuze_process.extractOne(text, g.titles.keys(), score_cutoff=96)
        if not res:
            g.max_custom_title
            out.put_title(text, g.max_custom_title + 1)
        else:
            notes.append(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
            out.put_title(res[0], g.titles[res[0]])
    else:
        out.put_title(text, 1)

    return notes

'''
这段代码定义了一个名为 `process_text_block` 的函数，用于处理文本块形状（shape）。
以下是代码的解释：
1. 定义了一个函数 `process_text_block`，接受一个文本块形状对象 `shape` 和一个未使用的参数 `_` 作为输入。
2. 声明 `out` 为全局变量。
3. 判断给定的文本块形状是否为列表块（调用 `is_list_block(shape)` 函数）。
   - 如果是列表块：
     - 遍历文本块形状中的每个段落（`shape.text_frame.paragraphs`）。
     - 如果当前段落的文本内容去除首尾空白字符后为空字符串，则继续下一次循环。
     - 获取格式化后的段落文本（调用 `get_formatted_text(para)` 函数）。
     - 调用 `out.put_list` 方法将格式化后的文本添加为列表项，同时传递段落的级别（`para.level`）。
     - 写入换行符到输出文件。
   - 如果不是列表块：
     - 遍历文本块形状中的每个段落（`shape.text_frame.paragraphs`）。
     - 如果当前段落的文本内容去除首尾空白字符后为空字符串，则继续下一次循环。
     - 获取格式化后的段落文本（调用 `get_formatted_text(para)` 函数）。
     - 调用 `out.put_para` 方法将格式化后的文本作为段落块输出。
4. 返回一个空列表 `[]`。
该函数用于处理给定的文本块形状，根据是否为列表块来决定如何处理文本内容。如果是列表块，遍历文本块中的每个段落，获取格式化后的段落文本，并将其作为列表项输出到输出文件中。
如果不是列表块，将文本块中的每个段落视为独立的段落块，并将格式化后的段落文本输出到输出文件中。最后，返回一个空列表表示没有备注信息。
'''
def process_text_block(shape, _):
    global out
    if is_list_block(shape):
        # generate list block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_formatted_text(para)
            out.put_list(text, para.level)
        out.write('\n')
    else:
        # generate paragraph block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_formatted_text(para)
            out.put_para(text)
    return []

'''
这段代码定义了一个名为 `process_picture` 的函数，用于处理图片形状（shape）。
以下是代码的解释：
1. 定义了一个函数 `process_picture`，接受一个图片形状对象 `shape` 和幻灯片索引 `slide_idx` 作为输入参数。
2. 声明 `notes` 列表，用于存储处理过程中的备注信息。
3. 如果全局变量 `g.disable_image` 为真，则直接返回空的备注信息列表，表示禁用图片处理。
4. 声明全局变量 `picture_count` 和 `out`。
5. 构造图片文件名，使用 `g.file_prefix` 和 `picture_count` 进行拼接。
6. 获取图片的文件扩展名（`shape.image.ext`）。
7. 如果输出图片的目录（`g.img_path`）不存在，则创建该目录。
8. 构造输出图片的完整路径（`output_path`）和相对路径（`img_outputter_path`）。
9. 使用二进制写入模式打开输出路径对应的文件，将图片的二进制数据写入文件中。
10. 增加 `picture_count` 的计数。
11. 如果图片扩展名不是 'wmf'（即普通图片）：
    - 调用 `out.put_image` 方法将图片输出到输出文件中，传递图片的相对路径和最大宽度（`g.max_img_width`）。
    - 返回空的备注信息列表。
12. 如果图片扩展名是 'wmf'（即 WMF 格式的图片）：
    - 尝试将 WMF 图片转换为 PNG 格式，如果转换成功，则使用 PNG 格式的图片输出到输出文件中，并将转换成功的备注信息添加到 `notes` 列表中。
    - 如果转换失败，将原始的 WMF 图片输出到输出文件中，并将转换失败的备注信息添加到 `notes` 列表中。
13. 返回 `notes` 列表，其中包含了处理过程中的备注信息。

该函数用于处理给定的图片形状，根据图片的类型（普通图片或 WMF 图片）进行相应的处理。对于普通图片，将图片输出到输出文件中。
对于 WMF 图片，尝试将其转换为 PNG 格式，如果转换成功，则使用 PNG 格式的图片输出；如果转换失败，则仍然使用原始的 WMF 图片输出。最后，返回处理过程中的备注信息列表。
'''
def process_picture(shape, slide_idx):
    notes = []
    if g.disable_image:
        return notes
    global picture_count
    global out

    pic_name = g.file_prefix + str(picture_count)
    pic_ext = shape.image.ext
    if not os.path.exists(g.img_path):
        os.makedirs(g.img_path)

    output_path = g.path_name_ext(g.img_path, pic_name, pic_ext)
    common_path = os.path.commonpath([g.out_path, g.img_path])
    img_outputter_path = os.path.relpath(output_path, common_path)
    with open(output_path, 'wb') as f:
        f.write(shape.image.blob)
        picture_count += 1

    # normal images
    if pic_ext != 'wmf':
        out.put_image(img_outputter_path, g.max_img_width)
        return notes

    # wmf images, try to convert, if failed, output as original
    try:
        Image.open(output_path).save(os.path.splitext(output_path)[0] + '.png')
        out.put_image(os.path.splitext(img_outputter_path)[0] + '.png', g.max_img_width)
        notes.append(f'Image {output_path} in slide {slide_idx} converted to png.')
    except Exception as e:
        notes.append(
            f'Cannot convert wmf image {output_path} in slide {slide_idx} to png, this probably won\'t be displayed correctly.'
        )
        out.put_image(img_outputter_path, g.max_img_width)
    return notes

'''
这段代码定义了一个名为 `process_table` 的函数，用于处理表格形状（shape）。
以下是代码的解释：
1. 定义了一个函数 `process_table`，接受一个表格形状对象 `shape` 和一个未使用的参数 `_` 作为输入。
2. 声明全局变量 `out`。
3. 通过列表推导式，获取表格中每个单元格的文本内容，并将其组织成二维列表形式的表格数据（`table`）。
4. 如果表格的行数大于 0，则调用 `out.put_table` 方法将表格数据输出到输出文件中。
5. 返回一个空列表 `[]`，表示没有备注信息。

该函数用于处理给定的表格形状，将表格中每个单元格的文本内容提取出来，并将其以二维列表的形式表示为表格数据。然后，将表格数据输出到输出文件中。最后，返回一个空列表表示没有备注信息。
'''
def process_table(shape, _):
    global out
    table = [[cell.text for cell in row.cells] for row in shape.table.rows]
    if len(table) > 0:
        out.put_table(table)
    return []

'''
这段代码定义了一个名为 `ungroup_shapes` 的函数，用于展开（解组）形状对象列表。
以下是代码的解释：
1. 定义了一个函数 `ungroup_shapes`，接受一个形状对象列表 `shapes` 作为输入参数。
2. 声明一个空列表 `res`，用于存储展开后的形状对象。
3. 遍历给定的形状对象列表 `shapes`：
   - 如果当前形状对象是一个组合形状（`shape.shape_type == MSO_SHAPE_TYPE.GROUP`）：
     - 递归调用 `ungroup_shapes` 函数，传递当前组合形状的子形状列表 `shape.shapes`。
     - 将递归调用的结果（展开后的形状对象）扩展到 `res` 列表中。
   - 否则，将当前形状对象添加到 `res` 列表中。
4. 返回展开后的形状对象列表 `res`。
该函数用于将给定的形状对象列表展开（解组）为一个扁平的形状对象列表。遍历形状对象列表，如果遇到组合形状，则递归地展开组合形状的子形状列表；
否则，直接将形状对象添加到结果列表中。最终返回展开后的形状对象列表。
'''
def ungroup_shapes(shapes):
    res = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            res.extend(ungroup_shapes(shape.shapes))
        else:
            res.append(shape)
    return res


'''
这段代码定义了一个名为 `parse` 的函数，用于解析（转换）PPTX文件。
以下是代码的解释：
1. 定义了一个函数 `parse`，接受一个 `prs` 对象和一个 `outputer` 对象作为输入参数。
2. 声明全局变量 `out`，将其赋值为 `outputer`，用于输出转换结果。
3. 声明一个空列表 `notes`，用于存储转换过程中的备注信息。
4. 使用 `enumerate` 函数遍历 `prs.slides`，获取幻灯片索引 `idx` 和幻灯片对象 `slide`。
5. 如果设置了特定的幻灯片页码（`g.page`），且当前幻灯片不是指定的页码，则跳过当前幻灯片的转换。
6. 声明一个空列表 `shapes`。
7. 尝试对当前幻灯片的形状对象进行展开和排序操作，将结果赋值给 `shapes` 列表。如果在展开和排序过程中出现异常，则打印错误信息，并打印出有问题的形状对象的类型、位置和尺寸信息。
8. 遍历当前幻灯片的形状对象列表 `shapes`：
   - 如果当前形状对象是标题形状（`is_title(shape)` 返回 `True`）：
     - 调用 `process_title` 函数处理标题形状，并将返回的备注信息添加到 `notes` 列表中。
   - 如果当前形状对象是文本块形状（`is_text_block(shape)` 返回 `True`）：
     - 调用 `process_text_block` 函数处理文本块形状，并将返回的备注信息添加到 `notes` 列表中。
   - 如果当前形状对象是图片形状（`shape.shape_type == MSO_SHAPE_TYPE.PICTURE`）：
     - 调用 `process_picture` 函数处理图片形状，并将返回的备注信息添加到 `notes` 列表中。
   - 如果当前形状对象是表格形状（`shape.shape_type == MSO_SHAPE_TYPE.TABLE`）：
     - 调用 `process_table` 函数处理表格形状，并将返回的备注信息添加到 `notes` 列表中。
9. 关闭输出器对象的输出文件。
10. 如果存在备注信息（`len(notes) > 0`）：
    - 打印 "Process finished with notice:" 的提示信息。
    - 遍历备注信息列表 `notes`，打印每条备注信息。
该函数用于解析给定的 PPTX 文件，将幻灯片中的标题、文本块、图片和表格等形状对象进行转换，并输出到指定的输出文件中。转换过程中的备注信息会被存储在 `notes` 
列表中，并在转换完成后打印出来。最后，返回转换过程中的备注信息列表。
'''
# main
def parse(prs, outputer):
    global out
    out = outputer
    notes = []
    for idx, slide in enumerate(tqdm(prs.slides, desc='Converting slides')):
        if g.page is not None and idx + 1 != g.page:
            continue
        shapes = []
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
        except:
            print('Bad shapes encountered in this slide. Please check or move them and try again.')
            print('shapes:')
            for sp in slide.shapes:
                print(sp.shape_type)
                print(sp.top, sp.left, sp.width, sp.height)

        for shape in shapes:
            if is_title(shape):
                notes += process_title(shape, idx + 1)
            elif is_text_block(shape):
                notes += process_text_block(shape, idx + 1)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                notes += process_picture(shape, idx + 1)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                notes += process_table(shape, idx + 1)
    out.close()

    if len(notes) > 0:
        print('Process finished with notice:')
        for note in notes:
            print(note)
