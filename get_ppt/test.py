# 导入python-pptx库
from pptx import Presentation
from datetime import datetime




def is_datetime(text):
    try:
        datetime.strptime(text, '%Y')
        return  True
    except ValueError:
        return  False


def contains_number_and_check_length(paragraph_text):
    # 将Paragraph对象转换为字符串
    # paragraph_text = str(string.text)

    # 检查字符串中是否包含数字
    contains_number = any(char.isdigit() for char in paragraph_text)

    # 检查数字是否大于字符串的长度，并且是否大于4
    if contains_number:
        if len(paragraph_text) <=4:
            return True
        number = int(''.join(filter(str.isdigit, paragraph_text)))
        is_greater_than_length = number > len(paragraph_text) <= 4
        return is_greater_than_length

    return False





# 打开PPT文件
ppt = Presentation(rf"F:\pptx2md\ppt\手捧绿叶保护环境PPT模板.pptx")
# 遍历每一张幻灯片
for slide in ppt.slides:
    # 获取幻灯片的编号
    slide_number = slide.slide_id
    # 打印幻灯片的编号
    print(f"幻灯片{slide_number}：")
    # 遍历幻灯片中的每一个形状
    for shape in slide.shapes:
        # 获取形状的类型
        shape_type = shape.shape_type
        # 判断是否是文本框对象
        if shape.has_text_frame:
            # 打印形状的类型
            #print(f"- 形状类型：{shape_type}")
            # 获取文本框对象
            text_frame = shape.text_frame
            # 获取文本内容
            text = text_frame.text.strip()
            text = ''.join(text.split())
            # if len(text) > 2:
            #     if is_datetime(text):
            #         continue
            #     elif '项目名' in text:
            #         continue
                # 打印文本内容
            print(f"- 文本内容：{text},{len(text)}")
        # 判断是否是GroupShape对象
        elif shape.shape_type == 6: # 6是GroupShape对象的类型编号
            # 获取GroupShape对象内部的形状列表
            subshapes = shape.shapes
            # 遍历形状列表，找到文本框对象
            for subshape in subshapes:
                # 判断是否是文本框对象
                if subshape.has_text_frame:
                    # 获取文本框对象
                    text_frame = subshape.text_frame
                    # 获取文本内容
                    text = text_frame.text.strip()
                    text =''.join(text.split())
                    # if len(text) > 2:
                    #     if is_datetime(text) or text=='LOGO':
                    #         continue
                    #     elif contains_number_and_check_length(text):
                    #          continue
                        # 打印文本内容
                    print(f"- GroupShape文本内容：{text},{len(text)}")
