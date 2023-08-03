from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_font_size(run):
    font_size = run.font.size
    if font_size is None:
        font_size = Pt(18)
    return font_size

def is_title(shape):
    for paragraph in shape.text_frame.paragraphs:
        # paragraph_text = paragraph.text
        # paragraph_text = ''.join(paragraph_text.split())
        for run in paragraph.runs:
            text = run.text.strip()
            return get_font_size(run) >= Pt(32)


def is_run_bold(run):
    if run.font.bold:
        return True
    return False

def is_all_english(input_string):
    # 移除空格和标点符号
    cleaned_string = "".join(char for char in input_string if char.isalpha())
    print(cleaned_string.isalpha())
    return cleaned_string.isalpha()


def is_word_title(shape):
    for paragraph in shape.text_frame.paragraphs:
        if '标题' in paragraph.text.strip():
            return True

def get_font_size_p(p):
    #获取段落的所有runs
    runs = p.runs
    font_sizes = set(run.font.size for run in runs )
    # 如果段落中所有的run都具有相同的字体大小，则返回该字体大小；否则返回None
    if len(font_sizes) == 1:
        return font_sizes.pop()
    else:
        return None

def counting_words(shape):
    if hasattr(shape, "text"):  # 检查是否有文本框
        text = shape.text  # 获取文本框内容
        count = len(text)  # 获取文本框字数
        if count != 0:  # 如果文本框为空，就跳过这个文本框
            width = shape.width  # 获取文本框宽度
            height = shape.height  # 获取文本框高度
            margin_left = shape.text_frame.margin_left  # 获取左边距
            margin_right = shape.text_frame.margin_right  # 获取右边距
            margin_top = shape.text_frame.margin_top  # 获取上边距
            margin_bottom = shape.text_frame.margin_bottom  # 获取下边距

            #print(width, height, margin_left, margin_right, margin_top, margin_bottom)
            #font_size = shape.text_frame.paragraphs[-1].font.size  # 获取字体大小
            #print(font_size)
            text_frame = shape.text_frame  # 获取文本框中的文本帧对象
            paragraph = text_frame.paragraphs[0]
            font_size = get_font_size_p(paragraph)
            if font_size is None:  # 如果字体大小没有值，就给它一个默认值
                font_size = 14
            else:
                font_size = font_size.pt
            line_spacing = shape.text_frame.paragraphs[0].line_spacing  # 获取行距
            if line_spacing is None:  # 如果行距没有值，就给它一个默认值
                line_spacing = 1.5

            area = (width - margin_left - margin_right) * (height - margin_top - margin_bottom)  # 计算有效面积
            # print("有效面积",area)
            char_area = font_size ** 2 * 0.3  # 计算每个字符的平均面积（假设中文字符）
            # print("占地面积",char_area*91400)
            #math.ceil()  向上取整
            chars_per_line = int((width - margin_left - margin_right) // (font_size * 0.3))  # 计算每行的字符数
            lines_per_page = int((height - margin_top - margin_bottom) // (font_size + line_spacing))  # 计算每页的行数
            max_chars = (chars_per_line) / 91400 * (lines_per_page / 91400) * 10  # 计算最多的字符数
            # print("文字:", text)
            # print("文字数量:", count)  # 打印文字数量
            # print("最多的字符数:", round(max_chars, 0))  # 打印最多的字符数（以字数为单位）
            # # print("比例:", count / max_chars) # 打印两者的比例
            # print("**************")
            return {"count_key": count,"max_chars_key":round(max_chars, 0)}



def get_textbox_text_and_position(slide,slide_index):
        textbox_data = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                    text = ' '.join(paragraphs)
                # text = shape.text
                    textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
            else:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                    text = ' '.join(paragraphs)
                    # text = shape.text
                    textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

        sorted_textbox_data = sorted(textbox_data, key=lambda entry: entry['left'])
        return sorted_textbox_data


def get_textbox_text_and_top(slide, slide_index):
    textbox_data = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                # text = shape.text
                if len(text) <=4:
                    continue
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
        else:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                if len(text) <=4:
                    continue
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

    sorted_textbox_data = sorted(textbox_data, key=lambda entry: entry['top'])
    # print(sorted_textbox_data)
    return sorted_textbox_data



def get_left_positions(slide):
    left_positions = []
    # 遍历幻灯片中的每个形状
    for shape in slide.shapes:
        # 判断形状是否是文本框
        if shape.has_text_frame:
            left_positions.append(shape.left)

    return left_positions





def is_text_bold(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if is_run_bold(run):
                return True
    return False

def contains_number_and_check_length(paragraph_text):
    # 将Paragraph对象转换为字符串
    # paragraph_text = str(string.text)
    # 检查字符串中是否包含数字
    contains_number = any(char.isdigit() for char in paragraph_text)
    # 检查数字是否大于字符串的长度，并且是否大于4
    if contains_number:
        if len(paragraph_text) <= 4:
            return True
        try:
            number = int(''.join(filter(str.isdigit, paragraph_text)))
            is_greater_than_length = number > len(paragraph_text) <= 4
            return is_greater_than_length
        except ValueError:
            # 转换失败，说明字符串中包含了非法的数字字符，返回False

            return True

    return False

import  re
def is_english_with_spaces(input_string):
    # 定义正则表达式，表示全是英文字符并且包含空格
    pattern = r'^[a-zA-Z\s]+$'

    # 使用 re.match() 方法尝试匹配正则表达式
    # 如果匹配成功，返回一个匹配对象，否则返回 None
    match = re.match(pattern, input_string)

    # 判断是否匹配成功，并返回结果
    return bool(match)


def read_presentation(shape):
    text_frame = shape.text_frame
    paragraphs = text_frame.paragraphs
    # 第一个段落通常被视为标题
    title = paragraphs[0].text.strip()
    # 从第二个段落开始为正文
    content = "".join([p.text.strip() for p in paragraphs[1:]] )
    for p in paragraphs[1:]:
        if len(title) == len(p.text.strip()):
            return False
    if len(title) != 0 and len(content) != 0:
        # print("Title:", title)
        # print("Content:", content)
        # print("---------------------")
        return [True,title,content]


def compare_font_sizes(slide):
    """
    获取 ppt 页面文本的字号大小并找到最大字号大小及其所属的字段值
    """
    max_font_size = 0
    max_font_text = ""
    # 遍历每个页面
    # for slide in prs.slides:
        #print(f"页面 {slide.slide_id}:")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            font_size_list = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
            # 如果段落中有多个运行（Run），你可能想根据具体需求处理这些字号大小
            if len(paragraph.text) > 0:
                # print(f"段落文本：'{paragraph.text}', 字号大小列表：{font_size_list}")
                max_paragraph_font_size = max(font_size_list, default=0)
                if max_paragraph_font_size > max_font_size:
                    max_font_size = max_paragraph_font_size
                    max_font_text = paragraph.text

    # print(f"页面最大字号大小：{max_font_size} pt")
    # print(f"对应最大字号大小的字段值：'{max_font_text}'")
    return [max_font_size, max_font_text]



# def read_ppt_content(ppt_file_path):
#     presentation = Presentation(ppt_file_path)
#     content_data = []
#     for slide_index, slide in enumerate(presentation.slides):
#         # if slide_index >= 3:
#         #     continue
#         left_ = get_textbox_text_and_position(slide,slide_index)
#         top_ = get_textbox_text_and_top(slide,slide_index)
#         #print(left_[0]['text'])
#         slide_content = {
#             "slide_index": slide_index,
#             "content": []
#         }
#         # 获取幻灯片上的所有形状，并按照它们的放置顺序排序
#         shapes = sorted(slide.shapes, key=lambda x: (x.left, x.top))
#         for shape in shapes:
#             # 获取形状的类型
#             shape_type = shape.shape_type
#             if shape.has_text_frame:
#                 c_w = counting_words(shape)
#                 if '前  言' in shape.text_frame.text or '前言' in shape.text_frame.text:
#                     continue
#                 elif '目录' in shape.text_frame.text or '目\n录' in shape.text_frame.text or '目 录' in shape.text_frame.text:
#                     continue
#                 # elif compare_font_sizes(slide)[0] >= 40:
#                 #     if contains_number_and_check_length(shape.text_frame.text):
#                 #        continue
#                 elif is_english_with_spaces(shape.text_frame.text):
#                       continue
#                 elif len(shape.text_frame.text) < 2:
#                     continue
#                 if shape.text_frame.text:
#                     ls = read_presentation(shape)
#                     text = shape.text_frame.text
#                     if ls:
#                             if ls[1]:
#                                 type_ = '标题'
#                                 item = {
#                                     'id':f'{shape.shape_id}',
#                                     'theme': '华尔街发展史',
#                                     "hint_ext": ls[1],
#                                     "type": type_,
#                                     "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                     "text": ''
#                                 }
#                                 slide_content["content"].append(item)
#                             if ls[2]:
#                                 type_ = '正文'
#                                 item = {
#                                     'id': f'{shape.shape_id}',
#                                     'theme': '华尔街发展史',
#                                     "hint_ext": ls[2],
#                                     "type": type_,
#                                     "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                     "text": ''
#                                 }
#                                 slide_content["content"].append(item)
#                     else:
#                             if slide_index == 0:
#                                 if compare_font_sizes(slide)[1] == text:
#                                     type_ = '封面标题'
#                                 # if (is_title(shape) or is_text_bold(shape) == True or is_word_title(shape)) and '副标题' not in shape.text_frame.text:
#                                 #     type_ = '封面标题'
#                                 # elif '副标题' in shape.text_frame.text:
#                                 #     type_ = '副标题'
#                                 else:
#                                     type_ = '副标题'
#                             elif slide_index ==1 or slide_index ==2:
#                                 if len(text) <= 14 and text.isdigit() == False:
#                                     print(text)
#                                     type_ = '目录标题'
#                                 else:
#                                     type_ = '正文'
#                             elif text.isdigit() and len(text) <= 2 :
#                                     # print(compare_font_sizes(slide)[0],compare_font_sizes(slide)[1])
#                                     type_ = '目录章节标题'
#                             elif left_[0]['text'] == text or top_[0]['text']==text and len(text) <=15:
#                                 type_ = '章节标题'
#                             elif '标题内容' in text and len(text) >20:
#                                 type_ = '正文'
#                             elif '副标题' in text:
#                                 type_ = '副标题'
#                             else:
#                                 if  is_title(shape) or is_text_bold(shape) == True or is_word_title(shape):
#                                         type_ = '标题'
#                                 else:
#                                         type_ = '正文'
#                             item = {
#                                 'id': f'{shape.shape_id}',
#                                 'theme': '华尔街发展史',
#                                 "hint_ext": text,
#                                 "type": type_,
#                                 "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                 "text": ''
#                             }
#                             slide_content["content"].append(item)
#             elif shape_type == 6:  # 6是GroupShape对象的类型编号
#                 # 获取GroupShape对象内部的形状列表
#                 subshapes = sorted(shape.shapes, key=lambda x: (x.left, x.top))
#                 #subshapes = shape.shapes
#                 # 遍历形状列表，找到文本框对象
#                 for subshape in subshapes:
#                     c_w = counting_words(subshape)
#                     # 判断是否是文本框对象
#                     if subshape.has_text_frame:
#                         # 获取文本框对象
#                         text_frame = subshape.text_frame
#                         # 获取文本内容
#                         text = text_frame.text.strip()
#                         text = ''.join(text.split())
#                         if '前  言' in text or '前言' in text:
#                             continue
#                         elif '目录' in text or '目\n录' in text:
#                             continue
#                         elif contains_number_and_check_length(text):
#                             continue
#                         elif is_english_with_spaces(text):
#                             continue
#                         ls = read_presentation(subshape)
#                         if ls:
#                             if ls[1]:
#                                 type_ = '标题'
#                                 item = {
#                                     'id': f'{subshape.shape_id}',
#                                     'theme': '华尔街发展史',
#                                     "hint_ext": ls[1],
#                                     "type": type_,
#                                     "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                     "text": ''
#                                 }
#                                 slide_content["content"].append(item)
#                             if ls[2]:
#                                 type_ = '正文'
#                                 item = {
#                                     'id': f'{subshape.shape_id}',
#                                     'theme': '华尔街发展史',
#                                     "hint_ext": ls[2],
#                                     "type": type_,
#                                     "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                     "text": ''
#                                 }
#                                 slide_content["content"].append(item)
#                         else:
#                             if text and len(text) > 2:
#                                 if len(text) > 20 or '内容' in text:
#                                     type_ = '正文'
#                                 else:
#                                     if slide_index == 0:
#                                         if is_title(subshape) or is_text_bold(subshape) == True or is_word_title(subshape):
#                                             type_ = '封面标题'
#                                         elif '副标题' in shape.text_frame.text:
#                                             type_ = '副标题'
#                                         else:
#                                             type_ = '正文'
#                                     elif slide_index == 1:
#                                         type_ = '目录标题'
#                                     elif is_title(subshape) or is_text_bold(subshape) == True or is_word_title(subshape):
#                                         type_ = '标题'
#                                     elif text.isdigit() and len(text) <= 2:
#                                         # print(compare_font_sizes(slide)[0],compare_font_sizes(slide)[1])
#                                         type_ = '目录章节标题'
#                                     else:
#                                         if slide_index == 1:
#                                             type_ = '标题'
#                                         else:
#                                             type_ = '正文'
#
#                                 item = {
#                                     'id': f'{subshape.shape_id}',
#                                     'theme': '华尔街发展史',
#                                     "hint_ext": f"{text}",
#                                     "type": type_,
#                                     "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
#                                     "text": ''
#                                 }
#
#                                 slide_content["content"].append(item)
#
#         content_data.append(slide_content)
#
#     return content_data



def read_ppt_content(ppt_file_path):
    presentation = Presentation(ppt_file_path)
    content_data = []

    def is_cover_slide(index):
        return index == 0

    def is_table_of_contents_slide(index):
        return index == 1 or index == 2

    def is_chapter_title_slide(text):
        return text.isdigit() and len(text) <= 2

    def is_section_title_slide(left_text, top_text):
        return left_text == text or top_text == text and len(text) <= 15

    def is_content_slide(text):
        return '标题内容' in text and len(text) > 20

    def is_subtitle_slide(text):
        return '副标题' in text

    def create_slide_item(shape_id, theme, type_, hint_ext, count_key, max_chars_key):
        return {
            'id': f'{shape_id}',
            'theme': theme,
            'hint_ext': hint_ext,
            'type': type_,
            'number of words': f"{count_key}" if max_chars_key == 0.0 else f"{max_chars_key}",
            'text': ''
        }

    def get_slide_type(slide_index, shape, text):
        text_length = len(text)
        if text_length != 0:
            if is_cover_slide(slide_index):
                return '封面标题' if compare_font_sizes(slide)[1] == text else '副标题'
            elif is_table_of_contents_slide(slide_index):
                return '目录标题' if text_length <= 14 and not text.isdigit() else '正文'
            elif is_chapter_title_slide(text):
                return '目录章节标题'
            elif is_section_title_slide(left_[0]['text'], top_[0]['text']):
                return '章节标题'
            elif is_content_slide(text):
                return '正文'
            elif is_subtitle_slide(text):
                return '副标题'
            else:
                return '标题' if is_title(shape) or is_text_bold(shape) or is_word_title(shape) else '正文'

    for slide_index, slide in enumerate(presentation.slides):
        left_ = get_textbox_text_and_position(slide, slide_index)
        top_ = get_textbox_text_and_top(slide, slide_index)

        slide_content = {
            "slide_index": slide_index,
            "content": []
        }

        shapes = sorted(slide.shapes, key=lambda x: (x.left, x.top))
        for shape in shapes:
            shape_type = shape.shape_type
            if shape.has_text_frame:
                ls = read_presentation(shape)
                c_w = counting_words(shape)
                text_frame = shape.text_frame
                # text = text_frame.text.strip()
                # text = ''.join(text.split())
                text = text_frame.text
                if len(text) < 1:
                    continue
                if any(keyword in text for keyword in ['前  言', '前言', '目录', '目\n录', '目 录','前 言']):
                    continue
                elif contains_number_and_check_length(text) or is_english_with_spaces(text):
                    continue
                if ls:
                    type_ = '标题'
                    hint_ext = ls[1]
                    item = create_slide_item(shape.shape_id, '华尔街发展史', type_, hint_ext, c_w['count_key'],
                                             c_w['max_chars_key'])
                    slide_content["content"].append(item)

                    type_ = '正文'
                    hint_ext = ls[2]
                    item = create_slide_item(shape.shape_id, '华尔街发展史', type_, hint_ext, c_w['count_key'],
                                              c_w['max_chars_key'])
                    slide_content["content"].append(item)
                else:
                    slide_type = get_slide_type(slide_index, shape, text)
                    item = {
                        'id': f'{shape.shape_id}',
                        'theme': '华尔街发展史',
                        "hint_ext": text,
                        "type": slide_type,
                        "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
                        "text": ''
                    }
                    slide_content["content"].append(item)

            elif shape_type == 6:  # GroupShape
                subshapes = sorted(shape.shapes, key=lambda x: (x.left, x.top))
                for subshape in subshapes:
                    if subshape.has_text_frame:
                        ls = read_presentation(subshape)
                        c_w = counting_words(subshape)
                        text_frame = subshape.text_frame
                        text = text_frame.text
                        # text = text_frame.text.strip()
                        # text = ''.join(text.split())
                        if len(text) < 1:
                            continue
                        if any(keyword in text for keyword in ['前  言', '前言', '目录', '目\n录','前 言']):
                            continue
                        elif contains_number_and_check_length(text) or is_english_with_spaces(text):
                            continue
                        if ls:
                            type_ = '标题'
                            hint_ext = ls[1]
                            item = create_slide_item(subshape.shape_id, '华尔街发展史', type_, hint_ext, c_w['count_key'],
                                                     c_w['max_chars_key'])
                            slide_content["content"].append(item)

                            type_ = '正文'
                            hint_ext = ls[2]
                            item = create_slide_item(subshape.shape_id, '华尔街发展史', type_, hint_ext, c_w['count_key'],
                                                     c_w['max_chars_key'])
                            slide_content["content"].append(item)
                        else:
                            slide_type = get_slide_type(slide_index, subshape, text)
                            item = {
                                'id': f'{subshape.shape_id}',
                                'theme': '华尔街发展史',
                                "hint_ext": text,
                                "type": slide_type,
                                "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
                                "text": ''
                            }
                            slide_content["content"].append(item)

        content_data.append(slide_content)

    return content_data


import json
if __name__ == "__main__":
    ppt_file_path = rf"F:\pptx2md\ppt\大气工作总结计划汇报PPT模板.pptx"
    content_data = read_ppt_content(ppt_file_path)
    # json_output = json.dumps(content_data, indent=2, ensure_ascii=False)
    # print(json_output)
    with open("F:\pptx2md\json\大气工作总结计划汇报PPT模板.json", "w", encoding="utf-8") as file:
        json.dump(content_data, file, indent=2, ensure_ascii=False)

    print("JSON文件已成功写入。")
