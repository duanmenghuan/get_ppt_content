from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def get_font_size(run):
    # 获取文本运行（run）对象的字体大小
    font_size = run.font.size
    # 检查字体大小是否未定义（为None）
    if font_size is None:
        # 如果字体大小未定义，则将其设置为默认字体大小18个点（Pt）
        font_size = Pt(18)
    # 返回字体大小，可能是原始字体大小或默认值
    return font_size


def is_title(shape):
    # 遍历形状对象中的每个段落
    for paragraph in shape.text_frame.paragraphs:
        # 遍历段落中的每个文本运行（run）
        for run in paragraph.runs:
            # 获取文本运行的文本内容，并移除首尾的空格
            text = run.text.strip()
            # 检查文本运行的字体大小是否大于等于32个点（Pt）
            # 如果是标题文本（字体大小大于等于32），则返回True
            return get_font_size(run) >= Pt(32)


def is_run_bold(run):
    # 检查文本运行（run）的字体是否设置为粗体
    if run.font.bold:
        # 如果文本运行的字体设置为粗体，返回True
        return True
    # 如果文本运行的字体没有设置为粗体，返回False
    return False


def is_all_english(input_string):
    # 移除空格和标点符号，只保留字符串中的字母（英文字母）
    cleaned_string = "".join(char for char in input_string if char.isalpha())
    # 打印清理后的字符串是否仅包含字母
    print(cleaned_string.isalpha())
    # 检查清理后的字符串是否全部由字母组成，是则返回True，否则返回False
    return cleaned_string.isalpha()


def is_word_title(shape):
    # 遍历形状对象中的每个段落
    for paragraph in shape.text_frame.paragraphs:
        # 移除段落中的首尾空格，并检查是否包含 '标题' 这个字符串
        if '标题' in paragraph.text.strip() and not '标题内容' in paragraph.text.strip():
            # 如果段落包含 '标题' 字符串，返回True
            return True


def get_font_size_p(p):
    # 获取段落的所有运行（run）
    runs = p.runs
    # 使用集合（set）存储所有运行的字体大小，因为集合不允许重复元素，所以会自动去除重复的字体大小
    font_sizes = set(run.font.size for run in runs)
    # 如果段落中所有的运行都具有相同的字体大小，则返回该字体大小；否则返回None
    if len(font_sizes) == 1:
        # 如果集合中只有一个元素，即所有运行的字体大小都相同，则返回该字体大小
        return font_sizes.pop()
    else:
        # 如果集合中有多个元素，即运行的字体大小不同，则返回None
        return None


def counting_words(shape):
    if hasattr(shape, "text"):  # 检查是否有文本框
        text = shape.text  # 获取文本框内容
        count = len(text)  # 获取文本框字数
        if count != 0:  # 如果文本框为空，就跳过这个文本框
            width = shape.width  # 获取文本框宽度
            height = shape.height  # 获取文本框高度
            margin_left = shape.text_frame.margin_left  # 获取左边距
            margin_right = shape.text_frame.margin_right  # 获取右边距
            margin_top = shape.text_frame.margin_top  # 获取上边距
            margin_bottom = shape.text_frame.margin_bottom  # 获取下边距

            #print(width, height, margin_left, margin_right, margin_top, margin_bottom)
            #font_size = shape.text_frame.paragraphs[-1].font.size  # 获取字体大小
            #print(font_size)
            text_frame = shape.text_frame  # 获取文本框中的文本帧对象
            paragraph = text_frame.paragraphs[0]
            font_size = get_font_size_p(paragraph)
            if font_size is None:  # 如果字体大小没有值，就给它一个默认值
                font_size = 14
            else:
                font_size = font_size.pt
            line_spacing = shape.text_frame.paragraphs[0].line_spacing  # 获取行距
            if line_spacing is None:  # 如果行距没有值，就给它一个默认值
                line_spacing = 1.5

            area = (width - margin_left - margin_right) * (height - margin_top - margin_bottom)  # 计算有效面积
            # print("有效面积",area)
            char_area = font_size ** 2 * 0.3  # 计算每个字符的平均面积（假设中文字符）
            # print("占地面积",char_area*91400)
            #math.ceil()  向上取整
            chars_per_line = int((width - margin_left - margin_right) // (font_size * 0.3))  # 计算每行的字符数
            lines_per_page = int((height - margin_top - margin_bottom) // (font_size + line_spacing))  # 计算每页的行数
            max_chars = (chars_per_line) / 91400 * (lines_per_page / 91400) * 10  # 计算最多的字符数
            # print("文字:", text)
            # print("文字数量:", count)  # 打印文字数量
            # print("最多的字符数:", round(max_chars, 0))  # 打印最多的字符数（以字数为单位）
            # # print("比例:", count / max_chars) # 打印两者的比例
            # print("**************")
            return {"count_key": count,"max_chars_key":round(max_chars, 0)}



def get_textbox_text_and_position(slide):
        textbox_data = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                    text = ' '.join(paragraphs)
                # text = shape.text
                    textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
            else:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                    text = ' '.join(paragraphs)
                    # text = shape.text
                    textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

        sorted_textbox_data = sorted(textbox_data, key=lambda entry: entry['left'])
        return sorted_textbox_data


def get_textbox_text_and_top(slide):
    textbox_data = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                # text = shape.text
                if len(text) <=4:
                    continue
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})
        else:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if shape.has_text_frame:
                text_frame = shape.text_frame
                paragraphs = [p.text.strip() for p in text_frame.paragraphs]
                text = ' '.join(paragraphs)
                if len(text) <=4:
                    continue
                textbox_data.append({"left": left, "top": top, "width": width, "height": height, "text": text})

    sorted_textbox_data = sorted(textbox_data, key=lambda entry: entry['top'])
    # print(sorted_textbox_data)
    return sorted_textbox_data



def get_left_positions(slide):
    left_positions = []
    # 遍历幻灯片中的每个形状
    for shape in slide.shapes:
        # 判断形状是否是文本框
        if shape.has_text_frame:
            left_positions.append(shape.left)

    return left_positions





def is_text_bold(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if is_run_bold(run):
                return True
    return False

def contains_number_and_check_length(shape,paragraph_text):
    max_font_size = 0
    max_font_text = ""
    for paragraph in shape.text_frame.paragraphs:
        font_size_list = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
        # 如果段落中有多个运行（Run），你可能想根据具体需求处理这些字号大小
        if len(paragraph.text) > 0:
            # print(f"段落文本：'{paragraph.text}', 字号大小列表：{font_size_list}")
            max_paragraph_font_size = max(font_size_list, default=0)
            if max_paragraph_font_size > max_font_size:
                max_font_size = max_paragraph_font_size
                max_font_text = paragraph.text
    if max_font_size >= 40 and len(max_font_text) <=2:
        return False
    else:
        # 检查字符串中是否包含数字
        contains_number = any(char.isdigit() for char in paragraph_text)
        # 检查数字是否大于字符串的长度，并且是否大于4
        if contains_number:
            if len(paragraph_text) <= 4:
                return True
            try:
                number = int(''.join(filter(str.isdigit, paragraph_text)))
                is_greater_than_length = number > len(paragraph_text) <= 4
                return is_greater_than_length
            except ValueError:
                # 转换失败，说明字符串中包含了非法的数字字符，返回False

                return True


import  re
def is_english_with_spaces(input_string):
    # 定义正则表达式，表示全是英文字符并且包含空格
    pattern = r'^[a-zA-Z\s]+$'

    # 使用 re.match() 方法尝试匹配正则表达式
    # 如果匹配成功，返回一个匹配对象，否则返回 None
    match = re.match(pattern, input_string)

    # 判断是否匹配成功，并返回结果
    return bool(match)

def is_english(text):
    for char in text:
        if char.isalpha() and not char.isascii():
            return False
    return True


def read_presentation(shape):
    text_frame = shape.text_frame
    paragraphs = text_frame.paragraphs
    # 第一个段落通常被视为标题
    title = paragraphs[0].text.strip()
    # 从第二个段落开始为正文
    content = "".join([p.text.strip() for p in paragraphs[1:]] )
    for p in paragraphs[1:]:
        if len(title) == len(p.text.strip()):
            return False
        elif is_english(title):
            return False
        elif len(title) < len(p.text.strip()):
            return [True, title, content]
        else:
            if p.alignment == PP_PARAGRAPH_ALIGNMENT.LEFT:
                return False
            else:
                return False

    if len(title) < len(content):
        # print("Title:", title)
        # print("Content:", content)
        # print("---------------------")
        return [True, title, content]


def is_directory_title(s):
    # 正则表达式模式，匹配 "PART" 后面可以是一个或两个数字，然后可以是空格，再后面可以是 "ONE" 到 "NINE" 或 "01" 到 "09"
    pattern = r"PART\s+(\d+|one|two|three|four|five|six|seven|eight|nine)"
    return re.search(pattern, s, re.IGNORECASE) is not None


def compare_font_sizes(slide):
    """
    获取 ppt 页面文本的字号大小并找到最大字号大小及其所属的字段值
    """
    max_font_size = 0
    max_font_text = ""
    # 遍历每个页面
    # for slide in prs.slides:
        #print(f"页面 {slide.slide_id}:")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            font_size_list = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
            # 如果段落中有多个运行（Run），你可能想根据具体需求处理这些字号大小
            if len(paragraph.text) > 0:
                # print(f"段落文本：'{paragraph.text}', 字号大小列表：{font_size_list}")
                max_paragraph_font_size = max(font_size_list, default=0)
                if max_paragraph_font_size > max_font_size:
                    max_font_size = max_paragraph_font_size
                    max_font_text = paragraph.text

    # print(f"页面最大字号大小：{max_font_size} pt")
    # print(f"对应最大字号大小的字段值：'{max_font_text}'")
    return [max_font_size, max_font_text]

def is_word(string):
    pattern = r"^[a-zA-Z]+$"
    match = re.match(pattern, string)
    return bool(match)



from langdetect import detect
def is_english_string(text):
    if len(text) >= 3:
        language = detect(text)
        if language == 'en':
            return True
        else:
            #print(f"The string is in {language} language, not an English string.")
            return False


def read_ppt_content(ppt_file_path):
    presentation = Presentation(ppt_file_path)
    content_data = []

    def is_cover_slide(index):
        return index == 0

    def is_table_of_contents_slide(index):
        return index == 1 or index == 2

    def is_chapter_title_slide(text):
        compare_font_size = compare_font_sizes(slide)
        #print(compare_font_size[0],compare_font_size[1],text)
        if compare_font_size[1].isdigit():
            print(compare_font_size[1],compare_font_size[0])
        return compare_font_size[1].isdigit()  and compare_font_size[0] >=40

    def is_section_title_slide(left_text, top_text):
        return (left_text == text or top_text == text) and len(text) <= 15

    def is_content_slide(text):
        if '标题内容' in text and len(text) > 20:
            return True
        elif is_english_string(text):
            return True
        elif '文本' in text:
            return  True

    def is_subtitle_slide(text):
        return '副标题' in text

    def create_slide_item(shape_id, theme, type_, hint_ext, count_key, max_chars_key):
        return {
            'id': f'{shape_id}',
            'theme': theme,
            'hint_ext': hint_ext,
            'type': type_,
            'number of words': f"{count_key}" if max_chars_key == 0.0 else f"{max_chars_key}",
            'text': ''
        }

    def get_slide_type(slide_index, shape, text):
        text_length = len(text)
       # print(slide_index,text_length,text)
        try:
            if text_length != 0:
                if is_cover_slide(slide_index):
                    return '封面标题' if compare_font_sizes(slide)[1] == text else ' '
                elif is_word(text) and (is_text_bold(shape)==False or is_title(shape)==False):
                    return ' '
                elif is_table_of_contents_slide(slide_index) and not is_directory_title(text)==True:
                    if is_chapter_title_slide(text):
                        return '目录章节标题'
                    else:
                         if text_length <= 18 and not text.isdigit():
                             return '目录标题'
                         elif text_length >=60 or is_content_slide(text):
                             return '正文'
                         else:
                             return " "
                elif is_chapter_title_slide(text) or is_directory_title(text):
                    return '目录章节标题'
                elif is_section_title_slide(left_[0]['text'], top_[0]['text']):
                    return '章节标题'
                elif is_content_slide(text):
                    return '正文'
                elif is_subtitle_slide(text):
                    return '副标题'
                # elif text.isalpha():
                #     return '正文'
                else:
                    return '标题' if is_title(shape) or is_text_bold(shape) or is_word_title(shape) else '正文'
        except  IndexError:
            #print("数组异常")
            # if is_chapter_title_slide(text):
            #         return '目录章节标题'
            # else:
            return '标题' if is_title(shape) or is_text_bold(shape) or is_word_title(shape) else '正文'



    for slide_index, slide in enumerate(presentation.slides):
        left_ = get_textbox_text_and_position(slide)
        top_ = get_textbox_text_and_top(slide)

        slide_content = {
            "slide_index": slide_index,
            "content": []
        }

        shapes = sorted(slide.shapes, key=lambda x: (x.left, x.top))
        for shape in shapes:
            shape_type = shape.shape_type
            if shape.has_text_frame:
                ls = read_presentation(shape)
                c_w = counting_words(shape)
                text_frame = shape.text_frame
                # text = text_frame.text.strip()
                # text = ''.join(text.split())
                text = text_frame.text
                # print(text)
                if not text.isdigit() and len(text) <=1:
                    continue
                if any(keyword in text for keyword in ['前  言', '前言', '目录', '目\n录', '目 录','前 言','MU  LU   '] if len(text.strip()) <=5):
                    continue
                elif contains_number_and_check_length(shape,text):
                    continue
                    # or is_english_with_spaces(text)
                if ls:
                    type_ = '标题'
                    hint_ext = ls[1]
                    item = create_slide_item(shape.shape_id, '日本动漫发展史', type_, hint_ext, c_w['count_key'],
                                             c_w['max_chars_key'])
                    slide_content["content"].append(item)

                    type_ = '正文'
                    hint_ext = ls[2]
                    item = create_slide_item(shape.shape_id, '日本动漫发展史', type_, hint_ext, c_w['count_key'],
                                              c_w['max_chars_key'])
                    slide_content["content"].append(item)
                else:
                    slide_type = get_slide_type(slide_index, shape, text)
                    item = {
                        'id': f'{shape.shape_id}',
                        'theme': '日本动漫发展史',
                        "hint_ext": text,
                        "type": slide_type,
                        "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
                        "text": ''
                    }
                    slide_content["content"].append(item)

            elif shape_type == 6:  # GroupShape
                subshapes = sorted(shape.shapes, key=lambda x: (x.left, x.top))
                for subshape in subshapes:
                    if subshape.has_text_frame:
                        ls = read_presentation(subshape)
                        c_w = counting_words(subshape)
                        text_frame = subshape.text_frame
                        text = text_frame.text
                        # text = text_frame.text.strip()
                        # text = ''.join(text.split())
                        if  not text.isdigit() and len(text) <=1:
                            continue
                        # print("GroupShape", text)
                        if any(keyword in text for keyword in ['前  言', '前言', '目录', '目\n录','前 言','MU  LU   '] if len(text.strip()) <=5):
                            continue
                        elif contains_number_and_check_length(subshape,text):
                             continue
                        # or is_english_with_spaces(text)
                        if ls:
                            type_ = '标题'
                            hint_ext = ls[1]
                            item = create_slide_item(subshape.shape_id, '日本动漫发展史', type_, hint_ext, c_w['count_key'],
                                                     c_w['max_chars_key'])
                            slide_content["content"].append(item)

                            type_ = '正文'
                            hint_ext = ls[2]
                            item = create_slide_item(subshape.shape_id, '日本动漫发展史', type_, hint_ext, c_w['count_key'],
                                                     c_w['max_chars_key'])
                            slide_content["content"].append(item)
                        else:
                            slide_type = get_slide_type(slide_index, subshape, text)
                            item = {
                                'id': f'{subshape.shape_id}',
                                'theme': '日本动漫发展史',
                                "hint_ext": text,
                                "type": slide_type,
                                "number of words": f"{c_w['count_key']}" if c_w['max_chars_key'] == 0.0 else f"{c_w['max_chars_key']}",
                                "text": ''
                            }
                            slide_content["content"].append(item)

        content_data.append(slide_content)

    return content_data


import json
if __name__ == "__main__":
    ppt_file_path = rf"F:\pptx2md\ppt\手捧绿叶保护环境PPT模板.pptx"
    content_data = read_ppt_content(ppt_file_path)
    # json_output = json.dumps(content_data, indent=2, ensure_ascii=False)
    # print(json_output)
    with open("F:\pptx2md\json\大气工作总结计划汇报PPT模板.json", "w", encoding="utf-8") as file:
        json.dump(content_data, file, indent=2, ensure_ascii=False)

    print("JSON文件已成功写入。")
