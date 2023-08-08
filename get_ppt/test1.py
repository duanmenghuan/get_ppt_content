from pptx import Presentation


def extract_text_from_shapes(ppt_file, target_page):
    prs = Presentation(ppt_file)

    shapes_text = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_num == target_page:
            for shape in slide.shapes:
                shape_type = shape.shape_type
                if hasattr(shape, 'text'):
                    shapes_text.append(shape.text)
                elif shape_type == 6:
                    if hasattr(shape, 'text'):
                        shapes_text.append(shape.text)

    return shapes_text


# 指定PPT文件路径和目标页数
ppt_file_path = rf"F:\pptx2md\ppt\手捧绿叶保护环境PPT模板.pptx"
target_page_number = 4  # 修改为您想要提取的页数

# 提取文本内容
extracted_text = extract_text_from_shapes(ppt_file_path, target_page_number)

# 打印提取的文本内容
for idx, text in enumerate(extracted_text, start=1):
    print(f"Shape {idx} Text:")
    print(text)
    print("-" * 30)
