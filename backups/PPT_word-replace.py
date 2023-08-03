#!/usr/bin/env python
# -*- coding:utf-8 -*-
# @Time  : 2023/7/24 11:17
# @Author: Jerry
# @File  : PPT_word-replace.py
import copy
from pptx import Presentation
from pptx.shapes.picture import Picture
import os
import sys
from pptx.enum.shapes import MSO_SHAPE_TYPE
import hashlib
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR, MSO_FILL
import time
import regex as re
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

time_start = time.time()
################txt部分代码
def read_input_txt(input_filepath):#symbol_word[md_slide][md_shape] = [md_word, line_num, change_check]
    symbol_word = {}
    # 假设input.txt中包含多行文本
    with open(input_filepath, "r", encoding='UTF-8') as input_md:
        lines = input_md.readlines()
    # 定义正则表达式
    pattern = r'^Slide_(\d+)_(\d+)_'
    # 遍历每一行文本，并进行匹配
    for line_num, line in enumerate(lines, start=0):
        match = re.match(pattern, line.strip())
        if match:
            # 提取Slide、Shape和word信息
            md_slide = int(match.group(1)) -1
            md_shape = int(match.group(2)) - 1
            md_word = line[match.end():].strip()
            change_check = None
            # 将信息保存到三维字典中
            if md_slide not in symbol_word:
                symbol_word[md_slide] = {}
            symbol_word[md_slide][md_shape] = [md_word, line_num, change_check]
    return symbol_word

#用于判断单个字符类型
def check_char_type(char):
    if re.match(r'[\p{IsHan}：，。、；？！“”（）【】{}]', char):
        return [char, 'character']
    elif re.match(r'[a-zA-Z]+|[\u0021-\u002F]+|[\u003A-\u0040]+|[\u005B-\u0060]+|[\u007B-\u007E]', char):
        return [char, 'letter']
    elif re.match(r'[0-9]', char):
        return [char, 'number']
    elif re.match(r'\u0020', char):
        return [char, 'space']
    else:
        return [char, 'other']

#用于判断是否应该拼接字符
def check_type_concatenate(category, current_category):
    if category == 'space':
        return True
    if category == current_category:
        return True
    if category == 'number' and current_category == 'letter':
        return True
    if category == 'letter' and current_category == 'number':
        return True
    else:
        return False

#用于判断并拼接字符
def merge_consecutive_data(data_dict):#没考虑空格开头的情况
    merged_data = []
    current_category = None
    current_text = ""
    for key, (char, category) in data_dict.items():
        if current_category == None:
            if current_text:
                merged_data.append((current_text, current_category))
            current_category = category
            current_text = char
        elif not check_type_concatenate(category, current_category):
            if current_text:
                merged_data.append((current_text, current_category))
            current_category = category
            current_text = char

        elif check_type_concatenate(category, current_category):
            current_text += char
        else:
            print('different situation!')
            print(current_category, category)
            print(current_text, char)
        if len(data_dict) == key + 1:
            if current_text:
                merged_data.append((current_text, current_category))
    return merged_data


#函数将拼接好的字符的标识符从 number letter转换为english,character转换为Chinese
#to do 细分:

def change_repl_key(replace_words):
    for line in replace_words:
        new_replace_words = []
        for words in replace_words[line]:
            if words[1] == 'number' or words[1] == 'letter':
                new_replace_words.append((words[0], 'english'))
            elif words[1] == 'character':
                new_replace_words.append((words[0], 'chinese'))
        replace_words[line] = new_replace_words
    return replace_words


#用于拼接同类字符
def if_rpl_words_chinese(replace_word):  #place__word = [line][place_num][char, family]
    result = {}
    place__word = {}
    for (word, line) in replace_word:
        place__word[line] = {}
        for place_num, char in enumerate(word, start=0):
            place__word[line][place_num] = {}
            place__word[line][place_num] = [check_char_type(char)[0], check_char_type(char)[1]]
    for line in place__word:
        result[line] = {}
        result[line] = merge_consecutive_data(place__word[line])
    return result

#用于读取replace文本to do:delete \n
def read_replace_txt(output_filepath):
    with open(output_filepath, encoding='UTF-8') as replace_md:
        replace__word = replace_md.readlines()
        replace_word = [(i.strip(), line_num) for line_num,i in enumerate(replace__word, start= 0)]
    replace___words = if_rpl_words_chinese(replace_word)#add language check
    replace_words = change_repl_key(replace___words)#chinese or english
    return replace_words


def save_run_font(run):
    save_run = {}
    font_size = None
    font_name = None
    if run.font.size is not None:
        font_size = run.font.size
    if run.font.name is not None:
        font_name = run.font.name
    if run.font.color.type == 1:
        color_type = 1
        color_value = run.font.color.rgb
    elif run.font.color.type == 2:
        color_type = 2
        color_value = run.font.color.theme_color
    # elif run.font.color.type == 'THEME_COLOR':
    #     color_type = 'THEME_COLOR'
    #     color_value = (run.funt.color.theme_color, prs.theme.theme_color_scheme.colors[run.funt.color.theme_color])#[index, value]
    else:
        color_type = None
        color_value = None
    save_run = {
        'text': run.text,
        'font_size': font_size,
        'font_name': font_name,
        'brightness': run.font.color.brightness,
        'bold': run.font.bold,
        'italic': run.font.italic,
        'underline': run.font.underline,
        'color_type': color_type,
        'color_value': color_value
    }
    return save_run


#读取现有run时遍历run内文本用于判断文本类型，添加标识符用于匹配
def if_runs_chinese(paragraph):
    eng_or_chn = None
    result_all = {}
    result_eng = {}
    result_chn = {}
    chn_num = []
    eng_num = []
    for run_num, run in enumerate(paragraph.runs, start=0):
        save_run = save_run_font(run)
        if hasattr(run, "text"):
            pattern_chinese = r'[\p{IsHan}：，。、；？！“”（）【】{}]'
            #pattern_punctuation = r'[\p{P}]'
            pattern_english = r'[a-zA-Z]+|[\u0021-\u002F]+|[\u003A-\u0040]+|[\u005B-\u0060]+|[\u007B-\u007E]'
            match_chinese = re.findall(pattern_chinese, run.text)
            #match_punctuation = re.findall(pattern_punctuation, run.text)
            match_english = re.findall(pattern_english, run.text)
            #拼接字符 防止纯字符文段
            #match_punctuation_all = ''.join(match_punctuation)
            if match_chinese:
                eng_or_chn = 'chinese'
                result_chn[run_num] = (save_run, False)
                chn_num.append(run_num)
            elif match_english:
                eng_or_chn = 'english'
                result_eng[run_num] = (save_run, False)
                eng_num.append(run_num)

            #elif match_punctuation_all == run.text:
            #    eng_or_chn = 'punctuation'
        result_all[run_num] = (save_run, eng_or_chn, False)
    print('lenth')
    print(len(result_chn))
    print(len(result_eng))
    print("result_all haven't been add")
    print('all:', result_all)
    print('eng:', result_eng)
    print('chn:', result_chn)
    print('chn_num:', chn_num)
    print('eng_num:', eng_num)
    #当不存在中文格式时，以最后一个英文格式作为中文格式
    if len(result_chn) == 0:
        print('chinese empty')
        print(len(result_chn))
        add_num = eng_num[len(eng_num) - 1] + 1
        chn_num.append(add_num)
        result_chn[add_num] = (result_eng[add_num - 1][0].copy(), False)
        print('chn:', result_chn[add_num]['font_name'])
        result_chn[add_num][0]['font_name'] = '微软雅黑'
        result_chn[add_num]['text'] = 'default chinese'
        print(f'add_default chinese run:<{chn_num[0]}>')
        result_all[add_num] = (result_chn[add_num][0], 'chinese', False)
        print(result_chn[add_num])
        print('lenth_chn after')
        print(len(result_eng))

    #当不存在英文格式时，以最后一个中文格式作为英文格式
    if len(result_eng) == 0:
        print('english empty')
        print(len(result_eng))
        add_num = chn_num[len(chn_num) - 1] + 1
        eng_num.append(add_num)
        result_eng[add_num] = (result_chn[add_num - 1][0].copy(), False)
        result_eng[add_num][0]['font_name'] = 'JetBrains Mono'
        result_eng[add_num][0]['text'] = 'default english'
        print(result_chn[add_num - 1][0]['text'])
        print(f'add_default english run:<{eng_num[0]}>')
        result_all[add_num] = (result_eng[add_num][0], 'english', False)

        print('lenth_eng after')
        print(len(result_eng))
    result = {
        'result_all': result_all,
        'result_chn': result_chn,
        'result_eng': result_eng,
        'chn_num': chn_num,
        'eng_num': eng_num
    }
    print('the result_all which is found')
    print('all:', result_all)
    print('eng:', result_eng)
    print('chn:', result_chn)
    return result


#用于处理多个run多个中英文混杂文本的情况时添加或减少新的run

def add_or_delete_runs(paragraph, count_runs):
    for i in range(count_runs):
        print(f'have <{len(paragraph.runs)}> runs')
        print(f'need <{count_runs}> runs')
        if len(paragraph.runs) == count_runs:
            break
        if len(paragraph.runs) < count_runs:
            new_run = paragraph.add_run()
            new_run.text = 'add_run'
            new_run.font.name = None
            print(f'now we have <{len(paragraph.runs)}> runs')

        if len(paragraph.runs) > count_runs:
            paragraph._p.remove(paragraph.runs[len(paragraph.runs)-1]._r)
    return paragraph



#对于run不可直接整体赋值 因为具有run和paragraph的对应关系 其中有部分不可直接赋值因此不能整体赋值单独写出一个函数，用于将run的所需的格式信息写到paragraph的run里
def clone_run_formatting(paragraph_run, saved_run):

    if saved_run['font_size'] is not None:
        paragraph_run.font.size = saved_run['font_size']
    print(saved_run['font_size'])
    if saved_run['font_name'] is not None:
        #paragraph_run.font.name = saved_run['font_name']
        paragraph_run.font.name = saved_run['font_name'][:]
    if saved_run['color_type'] == 1:
        paragraph_run.font.color.rgb = saved_run['color_value']
        paragraph_run.font.color.brightness = saved_run['brightness']
    # elif saved_run['color_type'] == "THEME_COLOR":
    #     paragraph_run.font.color.type = "THEME_COLOR"
    #     theme_color_index = saved_run['value'][0]
    #     theme_color_info = prs.theme.theme_color_scheme.colors[theme_color_index]
    #     paragraph_run.font.color.theme_color = theme_color_index
    #     paragraph_run.font.color.rgb = theme_color_info.rgb
    #     paragraph_run.font.color.brightness = saved_run['brightness']
    elif saved_run['color_type'] == 2:
        paragraph_run.font.color.theme_color = saved_run['color_value']
        paragraph_run.font.color.brightness = saved_run['brightness']

    paragraph_run.font.bold = saved_run['bold']
    paragraph_run.font.italic = saved_run['italic']
    paragraph_run.font.underline = saved_run['underline']





def clone_textbox_with_formatting(prs, Pos, run_info, replace_text):
    print('Pos：', Pos)
    slide = prs.slides[Pos[0]]
    shape = slide.shapes[Pos[1]]
    shape.text_frame.word_wrap = True

    paragraph_ = shape.text_frame.paragraphs[Pos[2]]
    paragraph = add_or_delete_runs(paragraph_, len(replace_text))
    print('replace_text:', replace_text)
    print(len(replace_text))
    print("目标结构：")
    for i in replace_text:
        print(i[1], end=' ')
    print("\n当前结构：")
    for i in run_info['result_all']:
        print(run_info['result_all'][i][1], end=' ')

    #添加格式标识flag
    eng_flag = 0
    chn_flag = 0
    for place in range(len(replace_text)):#place是正在读取的需要赋格式值的第place个文段
        # replace_text[place][0]是要换成的字 replace_text[place][1]是他的中英文识别符
        # result = {'result_all': result_all, 'result_chn': result_chn, 'result_eng': result_eng}
        # result_chn[run_num] = (run, False) result_eng[run_num] = (run, False)
        print('place:', place)
        if replace_text[place][1] == 'chinese':
            run_num = run_info['chn_num'][chn_flag]
            print('chn_run_num:', run_num)
            print('run_info_result_all used:', run_info['result_all'])
            print('run_info_result_eng used:', run_info['result_eng'])
            print('run_info_result_chn used:', run_info['result_chn'])
            print(run_info['eng_num'])
            print(run_info['chn_num'])
            run_info_ = (run_info['result_chn'][run_num])
            clone_run_formatting(paragraph.runs[place], run_info_[0])
            paragraph.runs[place].text = replace_text[place][0]
            print(f'当前位置:{place} 使用的是原paragraph的: runs[{run_num}] 的格式')
            print(f"它的文本是：<{replace_text[place][0]}>  原始run的文本是：{run_info_[0]['text']}")
            print(f'原始的颜色type:', {run_info_[0]['color_type']})
            print(f'新的颜色type:', {paragraph.runs[place].font.color.type})
            print(f'paragraph.font.name: {paragraph.font.name}')
            print(f"原始字体: {run_info_[0]['font_name']}")
            print(f'当前字体： {paragraph.runs[place].font.name}')
            print(len(paragraph.runs))
            print('result_all:', run_info['result_all'])
            if chn_flag < len(run_info['chn_num']) - 1:
                chn_flag = chn_flag + 1
            print(len(run_info['chn_num']))
            print(chn_flag)
            print('________________________________________')
        if replace_text[place][1] == 'english':
            run_num = run_info['eng_num'][eng_flag]
            print('eng_run_num:', run_num)
            print('run_info_result_all used:', run_info['result_all'])
            print('run_info_result_eng used:', run_info['result_eng'])
            print('run_info_result_chn used:', run_info['result_chn'])
            print(run_info['eng_num'])
            print(run_info['chn_num'])
            run_info__ = (run_info['result_eng'][run_num])
            clone_run_formatting(paragraph.runs[place], run_info__[0])
            paragraph.runs[place].text = replace_text[place][0]
            print(f'当前位置{place} 使用的是原paragraph的: runs[{run_num}] 的格式')
            print(f"它的文本是：<{paragraph.runs[place].text}>  原始run的文本是：{run_info__[0]['text']}")
            print(f'原始的颜色type:', {run_info__[0]['color_type']})
            print(f'新的颜色type:', {paragraph.runs[place].font.color.type})
            print(f'paragraph.font.name: {paragraph.font.name}')
            print(f"原始字体: {run_info__[0]['font_name']}")
            print(f'当前字体{paragraph.runs[place].font.name}')
            print(len(paragraph.runs))
            print('result_all:', run_info['result_all'])
            if eng_flag < len(run_info['eng_num']) - 1:
                eng_flag = eng_flag + 1
            print(len(run_info['eng_num']))
            print(eng_flag)
            print('________________________________________')

    # left = shape.left
    # top = shape.top
    # width = shape.width
    # height = shape.height
    # 创建新的文本框
    #shape.text_frame.word_wrap = True
    # 复制格式信息到新的文本
    #print(text_info["replace_text"])



################替换文本
def ReplaceSlideText(prs, slide_num, symbol_shapes, repl_words):
    for shape_num, shape in enumerate(prs.slides[slide_num].shapes, start=0):
        # left = shape.left
        # top = shape.top
        # width = shape.width
        # height = shape.height
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph_num, paragraph in enumerate(text_frame.paragraphs, start=0):
                text = paragraph.text.strip()
                if not text:
                    continue

                for shape_key, value in symbol_shapes.items():
                    runs = []
                    text_info = []
                    for run_num, run in enumerate(paragraph.runs, start=0):
                        # print(slide_num)
                        # print(shape_key)
                        # print(f'run.text:{run.text}')
                        # print(f'value:{value}')
                        # print(f'paragraph.text:{paragraph.text}')

                        #######获取格式信息
                        color_type = None
                        color_value = None
                        brightness = None
                        font_size = None
                        if run.font.size is not None:
                            font_size = run.font.size.pt
                        if paragraph.runs:
                            color = run.font.color
                            if color.type == MSO_COLOR_TYPE.RGB:
                                color_type = 'RGB'
                                color_value = color.rgb
                                brightness = color.brightness
                            elif color.type == MSO_COLOR_TYPE.SCHEME:
                                color_type = 'SCHEME'
                                color_value = color.theme_color
                                brightness = color.brightness
                            elif color.type == MSO_THEME_COLOR:
                                theme_color_index = color.theme_color
                                theme_color_info = prs.theme.theme_color_scheme.colors[theme_color_index]
                                print(theme_color_info)
                        else:
                            continue
                        ############################run内匹配
                        if value[0] == run.text and shape_num == shape_key and value[2] == None:
                        #to do:check shape
                            cur_text = run.text
                            #print('check',repl_words)
                            replace_text = cur_text.replace(value[0], (repl_words[value[1]][0][0]))
                            run.text = replace_text

                            TestSet = f'  TestSet:\n    font_name: {run.font.name}  color_type:{color_type}  color_value:{color_value}  brightness:{brightness}'
                            if run.font.bold:
                                TestSet += f'  bold:{run.font.bold}'
                            if run.font.italic:
                                TestSet += f'  italic:{run.font.italic}'
                            if run.font.underline:
                                TestSet += f'  underline:{run.font.underline}'
                            #print(f'Slide:{slide_num} Shape:{shape_num} changing:\nfrom <{cur_text}> to <{replace_text}>')
                            #print(f'  TextPos:\n    left:{left+(width/2)} top:{top+(height/2)}\n    width:{shape.width} height:{shape.height}\n    font_size:{font_size}')
                            #print(TestSet)
                            value[2] = True
                        #################paragraph匹配
                        elif value[0] != run.text and value[0] == text and shape_num == shape_key and value[2] == None:
                            run_info = if_runs_chinese(paragraph)
                            Pos = (slide_num, shape_num, paragraph_num)
                            replace_text = repl_words[value[1]]#replace_text是一行数据，包括文本和中英文标识符
                            clone_textbox_with_formatting(prs, Pos, run_info, replace_text)
                            value[2] = True




currpwd = os.path.dirname(os.path.abspath(__file__))
filepath = f'{currpwd}/fixtures/test001.pptx'
PicResavePath = f'{currpwd}/pic_save'
PicReplacePath = f'{currpwd}/pic_replace'
symbol_input_path = f'{currpwd}/input.txt'
replace_input_path = f'{currpwd}/replace.txt'


#读取ppt
prs = Presentation(filepath)
#读取txt
symbol_word = read_input_txt(symbol_input_path)
print(symbol_word)
replace_words = read_replace_txt(replace_input_path)
print(replace_words)

# 遍历ppt
for i in range(len(prs.slides)):
    print('the slide which is generating:', i)
    # 替换ppt文本
    # symbol_word识别词 用于查找md内对应部分
    symbol_shapes = symbol_word[i]
    ReplaceSlideText(prs, i, symbol_shapes, replace_words)

    # replace_word替换词 用于填充PPT内容
# 保存ppt

new_runs_order = [prs.slides[3].shapes[8].text_frame.paragraphs[0].runs[1], prs.slides[3].shapes[8].text_frame.paragraphs[0].runs[0]]
paragraph = prs.slides[3].shapes[8].text_frame.add_paragraph()
run1 = paragraph.add_run()
run1.text = "This is the first run. "
run1.font.size = Pt(14)

run2 = paragraph.add_run()
run2.text = "This is the second run. "
run2.font.size = Pt(18)

run3 = paragraph.add_run()
run3.text = "This is the third run. "
run3.font.size = Pt(12)


re_save_path = f'{currpwd}/ppt_outputs/re-place_output.pptx'
prs.save(re_save_path)

time_end = time.time()
print('time cost', time_end - time_start, 's')